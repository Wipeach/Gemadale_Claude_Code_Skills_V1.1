# -*- coding: utf-8 -*-
"""
Gemdale Slide Master Generator (updated)
Creates a PowerPoint presentation template with consistent branding
Changed: default slides reduced to 5 and page titles updated per user's request.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from PIL import Image
import os
from datetime import datetime
from pathlib import Path
from typing import Dict, Any


def create_slide_with_header_footer(prs, slide, slide_num, total_slides, header_image_path=None):
    """
    Add consistent header and footer to an individual slide
    """
    
    # Create header area
    header_height = Inches(1.2)
    
    if header_image_path and os.path.exists(header_image_path):
        try:
            header_img = Image.open(header_image_path)
            img_width, img_height = header_img.size
            aspect_ratio = img_height / img_width
            
            # Calculate dimensions to fit slide width
            image_width = prs.slide_width
            image_height = min(image_width * aspect_ratio, header_height)
            
            # Add header image
            pic = slide.shapes.add_picture(header_image_path, Inches(0), Inches(0), 
                                         width=image_width, height=image_height)
        except Exception as e:
            print(f"Warning: Could not add header image: {e}")
            # Create placeholder rectangle if image fails
            header_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
                prs.slide_width, header_height
            )
            header_shape.fill.solid()
            header_shape.fill.fore_color.rgb = RGBColor(0, 102, 204)  # Blue
            header_shape.line.fill.background()
    else:
        # Create placeholder header rectangle
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), 
            prs.slide_width, header_height
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(0, 102, 204)  # Blue
        header_shape.line.fill.background()
    
    # Create footer section with square blank space on left
    footer_height = Inches(0.8 * 0.8)  # 20% reduction
    footer_top = prs.slide_height - footer_height
    square_size = footer_height  # Same as footer height for square blank
    
    # Main footer rectangle (not full width, starts after square blank)
    footer_left = square_size  # Start after the blank square space
    footer_width = prs.slide_width - square_size
    
    footer_shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        footer_left, 
        footer_top, 
        footer_width, 
        footer_height
    )
    footer_shape.fill.solid()
    footer_shape.fill.fore_color.rgb = RGBColor(246, 75, 48)  # Changed to RGB(246,75,48)
    footer_shape.line.fill.background()
    
    # Colored square that "jumps up" above the blank area
    jump_square_size = Inches(0.64)  # Smaller jumping square
    jump_square_left = Inches(0)  # Positioned within the blank square area
    jump_square_bottom = footer_top - Inches(0.64)  # Slightly above footer
    
    jump_square = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        jump_square_left,
        jump_square_bottom,
        jump_square_size,
        jump_square_size
    )
    jump_square.fill.solid()
    jump_square.fill.fore_color.rgb = RGBColor(246, 75, 48)  # Gold/yellow for dynamic effect
    jump_square.line.fill.background()
    
    # Page number positioned in the blank square area (white background)
    page_num_left = Inches(0.05)  # Very close to left edge within blank area
    page_num_top = footer_top + Inches(0.15) # Align with center
    page_num_width = square_size - Inches(0.1)  # Fit within square blank
    page_num_height = footer_height
    
    page_num_shape = slide.shapes.add_textbox(
        page_num_left, page_num_top, page_num_width, page_num_height
    )
    
    page_num_tf = page_num_shape.text_frame
    page_num_tf.text = f"{slide_num}"  # Actual page number
    page_num_para = page_num_tf.paragraphs[0]
    page_num_para.alignment = PP_ALIGN.CENTER
    page_num_para.font.size = Pt(12)  # Slightly larger for visibility
    page_num_para.font.color.rgb = RGBColor(246, 75, 48)
    page_num_para.font.name = 'Arial'
    page_num_para.font.bold = True
    
    # Company name
    company_text_left = page_num_left + page_num_width + Inches(0.5)
    company_width = Inches(3)
    
    company_shape = slide.shapes.add_textbox(
        company_text_left, page_num_top, company_width, page_num_height
    )
    
    company_tf = company_shape.text_frame
    company_tf.text = "Gemdale Corporation"
    company_para = company_tf.paragraphs[0]
    company_para.alignment = PP_ALIGN.LEFT
    company_para.font.size = Pt(10)
    company_para.font.color.rgb = RGBColor(255, 255, 255)
    company_para.font.name = 'Arial'
    
    # Date and copyright info on right
    copyright_left = prs.slide_width - Inches(4.0)
    copyright_width = Inches(3.5)
    
    copyright_shape = slide.shapes.add_textbox(
        copyright_left, page_num_top, copyright_width, page_num_height
    )
    
    copyright_tf = copyright_shape.text_frame
    current_year = datetime.now().year
    copyright_tf.text = f"© {current_year} Gemdale Corp. - All Rights Reserved"
    copyright_para = copyright_tf.paragraphs[0]
    copyright_para.alignment = PP_ALIGN.RIGHT
    copyright_para.font.size = Pt(9)
    copyright_para.font.color.rgb = RGBColor(200, 200, 200)  # Light gray
    copyright_para.font.name = 'Arial'


def create_gemdale_slide_master_template():
    """
    Create a PowerPoint presentation template with consistent branding
    """
    
    # Create a new presentation
    prs = Presentation()
    
    # Set 16:9 aspect ratio (13.33" x 7.5")
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    return prs


def create_slides(prs, project_name: str, num_slides: int = 5, header_image_path: str = None):
    """
    Create slides with consistent header, footer, and page numbering
    Now accepts project_name so titles can include the project name.
    Default number of slides set to 5 and titles updated.
    """
    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]
    
    # Predefined titles (per user's request)
    predefined_titles = [
        f"x.x 一级竞品：{project_name}-基础信息",
        f"x.x 一级竞品：{project_name}-分批推售情况",
        f"x.x 一级竞品：{project_name}-户型分析",
        f"x.x 一级竞品：{project_name}-配置分析",
        f"x.x 一级竞品：{project_name}-客户分析",
    ]
    
    # Create slides
    for i in range(0, num_slides):
        slide = prs.slides.add_slide(content_layout)
        slide_num = i + 1  # Start numbering from 1 for the first slide
        
        # Remove the default title placeholder and content placeholder if present
        try:
            title_placeholder = slide.shapes.title
            if title_placeholder:
                title_placeholder.element.getparent().remove(title_placeholder.element)
        except Exception:
            pass
        
        try:
            # Many slide layouts have placeholder index 1 for content
            content_placeholder = slide.placeholders[1]
            if content_placeholder:
                content_placeholder.element.getparent().remove(content_placeholder.element)
        except Exception:
            pass
        
        # Add header and footer elements first
        create_slide_with_header_footer(prs, slide, slide_num, num_slides, header_image_path)
        
        # Create title text box that appears below the header with dynamic sizing
        header_height = Inches(1.2)
        square_size = Inches(0.8 * 0.8)
        offset_multiplier = 1.5
        title_left = square_size * offset_multiplier  # 1.5x of square size offset from left
        
        # Create title text box with appropriate initial dimensions for horizontal formatting
        title_left_margin = title_left  # Already calculated as 1.5x square size from left
        title_top_margin = Inches(0.1)
        title_width = prs.slide_width - title_left_margin - Inches(1.0)  # Leave right margin
        title_height = Inches(0.4)  # Reasonable height for a title
        
        title_shape = slide.shapes.add_textbox(title_left_margin, title_top_margin, title_width, title_height)
        title_tf = title_shape.text_frame
        title_tf.word_wrap = True
        
        # Choose the predefined title if available, otherwise fallback
        if i < len(predefined_titles):
            title_str = predefined_titles[i]
        else:
            title_str = f"标题：第{i+1}页标题"
        
        title_tf.text = title_str
        
        # Set text formatting
        title_para = title_tf.paragraphs[0]
        title_para.alignment = PP_ALIGN.LEFT
        title_para.font.name = '黑体'
        title_para.font.size = Pt(24)
        title_para.font.color.rgb = RGBColor(0, 0, 0)  # Dark text for white background
        title_para.font.bold = True
        
        # Ensure text fits properly in the frame
        title_shape.text_frame.auto_size = True
    
    return prs


def run(project_name: str, header_image_path: str = None, num_slides: int = 5, output_file: str = None) -> Dict[str, Any]:
    """
    Run the slide master creation with a given project name and optional parameters.
    
    Parameters:
    - project_name: str, name of the project
    - header_image_path: str, path to header image (optional)
    - num_slides: int, number of slides to create (default: 5)
    - output_file: str, output PPTX file path (optional)
    
    Returns:
    - dict: Contains the output file path and slide count
    """
    try:
        # Create the base presentation
        prs = create_gemdale_slide_master_template()
        
        # Check for header image
        if header_image_path is None:
            header_image_path = "./resources/images/gemdale_header.png"
        if os.path.exists(header_image_path):
            print(f"[OK] Header image found: {header_image_path}")
        else:
            print(f"[INFO] No header image found, using placeholder: {header_image_path}")
            header_image_path = None
        
        # Create slides with consistent branding
        prs = create_slides(prs, project_name=project_name, num_slides=num_slides, header_image_path=header_image_path)
        
        # Set output file path
        timestamp = datetime.now().strftime("%Y%m%d")
        if output_file is None:
            output_dir = Path(f"resources/working_data/{project_name}_{timestamp}/processed_data")
            output_dir.mkdir(parents=True, exist_ok=True)
            output_file = output_dir / f"{project_name}_gemdale_housing_project_template.pptx"
        else:
            output_file = Path(output_file)
            output_file.parent.mkdir(parents=True, exist_ok=True)
        
        # Save the presentation
        prs.save(output_file)
        
        print("[OK] Successfully created presentation template!")
        print(f"[FILE] {output_file}")
        size_width = prs.slide_width // 914400
        size_height = prs.slide_height // 914400
        print(f"[SIZE] {size_width}x{size_height} inches (16:9)")
        print(f"[SLIDES] {len(prs.slides)}")
        
        return {
            "output_file": str(output_file),
            "slide_count": len(prs.slides),
            "slide_size": f"{size_width}x{size_height} inches"
        }
    
    except Exception as e:
        print(f"[ERROR] Error creating presentation: {e}")
        import traceback
        traceback.print_exc()
        return {"output_file": "", "slide_count": 0, "slide_size": "", "error": str(e)}


if __name__ == "__main__":
    # For testing purposes, use a default project name (5 slides by default)
    result = run(project_name="华发四季半岛")
    print("\nSlide master creation result:", result)
