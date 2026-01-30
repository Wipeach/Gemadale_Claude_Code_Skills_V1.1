#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Insert multiple images into a PPT slide (page 3) with centered titles,
automatically arranged based on slide width.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, List

EMU_PER_INCH = 914400

def _compute_left_top_by_anchor(left_in: float, top_in: float, w_in: float, h_in: float, anchor: str) -> tuple[float, float]:
    """Convert anchor point to top-left coordinates."""
    anchor = (anchor or "top_left").lower()
    if anchor == "top_left":
        return left_in, top_in
    elif anchor == "top_right":
        return left_in - w_in, top_in
    elif anchor == "center":
        return left_in - w_in / 2.0, top_in - h_in / 2.0
    elif anchor == "bottom_left":
        return left_in, top_in - h_in
    elif anchor == "bottom_right":
        return left_in - w_in, top_in - h_in
    else:
        return left_in, top_in

def insert_images_with_titles(
    pptx_file_path: str,
    image_paths: List[str],
    slide_number: int = 3,
    top_in: float = 4.5,
    img_w_in: float = 1.8,
    side_margin_in: float = 0.5,
    title_height_in: float = 0.35,
    title_margin_in: float = 0.08,
    title_font_name: str = "Arial",
    title_font_size_pt: int = 12,
    auto_add_slide_if_missing: bool = True,
    output_suffix: str = ""
) -> bool:
    """
    Insert multiple images into a PPT slide with centered titles.
    
    Parameters:
    - pptx_file_path: str, path to PowerPoint file
    - image_paths: List[str], list of image file paths
    - slide_number: int, slide number (1-based indexing)
    - top_in: float, top position of images in inches
    - img_w_in: float, width of each image in inches
    - side_margin_in: float, margin on each side of slide in inches
    - title_height_in: float, height of title textbox in inches
    - title_margin_in: float, margin between title and image in inches
    - title_font_name: str, font name for titles
    - title_font_size_pt: int, font size for titles in points
    - auto_add_slide_if_missing: bool, whether to add blank slide if needed
    - output_suffix: str, unused since output is same as input
    
    Returns:
    - bool: True if successful, False otherwise
    """
    try:
        if not os.path.exists(pptx_file_path):
            print(f"[ERROR] PPTX file not found: {pptx_file_path}")
            return False

        prs = Presentation(pptx_file_path)

        if slide_number < 1:
            print(f"[ERROR] Invalid slide number: {slide_number}")
            return False
        while auto_add_slide_if_missing and slide_number > len(prs.slides):
            blank_layout = prs.slide_layouts[6]
            prs.slides.add_slide(blank_layout)

        if slide_number > len(prs.slides):
            print(f"[ERROR] Invalid slide number (exceeds slide count): {slide_number}")
            return False

        slide = prs.slides[slide_number - 1]

        imgs = [p for p in image_paths if os.path.exists(p)]
        count = len(imgs)

        if count == 0:
            print("[WARNING] No valid images found, skipping insertion.")
            prs.save(pptx_file_path)
            print(f"[INFO] Saved PPT without images to: {pptx_file_path}")
            return True

        page_w_in = prs.slide_width / EMU_PER_INCH
        usable_w_in = page_w_in - 2 * side_margin_in
        cell_w_in = usable_w_in / count
        real_img_w_in = min(img_w_in, cell_w_in * 0.95)

        for idx, img_path in enumerate(imgs):
            with Image.open(img_path) as im:
                w_px, h_px = im.size
            if w_px <= 0 or h_px <= 0:
                print(f"[ERROR] Invalid image size for {img_path}")
                continue

            aspect = h_px / float(w_px)
            real_img_h_in = real_img_w_in * aspect

            cell_left_in = side_margin_in + idx * cell_w_in
            left_in = cell_left_in + (cell_w_in - real_img_w_in) / 2
            tl_left_in, tl_top_in = _compute_left_top_by_anchor(left_in, top_in, real_img_w_in, real_img_h_in, "top_left")

            pic = slide.shapes.add_picture(
                img_path,
                Inches(tl_left_in),
                Inches(tl_top_in),
                Inches(real_img_w_in),
                Inches(real_img_h_in)
            )

            title_top_emu = pic.top - Inches(title_height_in) - Inches(title_margin_in)
            if title_top_emu < Inches(0.1):
                title_top_emu = Inches(0.1)

            textbox = slide.shapes.add_textbox(
                pic.left,
                title_top_emu,
                pic.width,
                Inches(title_height_in)
            )
            tf = textbox.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            p.text = f"户型 {idx+1}"
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.name = title_font_name
            run.font.size = Pt(title_font_size_pt)
            run.font.bold = True

            print(f"[INFO] Inserted: {img_path} | left={tl_left_in:.2f}in, top={tl_top_in:.2f}in, w={real_img_w_in:.2f}in, h={real_img_h_in:.2f}in")

        prs.save(pptx_file_path)
        print(f"[SUCCESS] Saved PPT with images to: {pptx_file_path}")
        return True

    except Exception as e:
        print(f"[ERROR] Failed to insert images: {str(e)}")
        return False

def run(
    project_name: str,
    pptx_file_path: str = None,
    image_paths: List[str] = None,
    slide_number: int = 3,
    top_in: float = 4.5,
    img_w_in: float = 1.8,
    side_margin_in: float = 0.5,
    title_height_in: float = 0.35,
    title_margin_in: float = 0.08,
    title_font_name: str = "Arial",
    title_font_size_pt: int = 12,
    auto_add_slide_if_missing: bool = True
) -> Dict[str, Any]:
    """
    Run the image insertion with titles on page 3 for a given project name.
    
    Parameters:
    - project_name: str, name of the project
    - pptx_file_path: str, path to PowerPoint file (optional, defaults to specified template)
    - image_paths: List[str], list of image file paths (optional, defaults to room_style*.jpg)
    - slide_number: int, slide number (1-based indexing)
    - top_in: float, top position of images in inches
    - img_w_in: float, width of each image in inches
    - side_margin_in: float, margin on each side of slide in inches
    - title_height_in: float, height of title textbox in inches
    - title_margin_in: float, margin between title and image in inches
    - title_font_name: str, font name for titles
    - title_font_size_pt: int, font size for titles in points
    - auto_add_slide_if_missing: bool, whether to add blank slide if needed
    
    Returns:
    - dict: Contains success status, output file path, and image paths
    """
    try:
        timestamp = datetime.now().strftime("%Y%m%d")
        if pptx_file_path is None:
            pptx_file_path = str(Path(f"resources/working_data/{project_name}_{timestamp}/processed_data/{project_name}_gemdale_housing_project_template.pptx"))
        if image_paths is None:
            image_paths = [
                str(Path(f"resources/images/room_style{i}.jpg")) for i in range(1, 6)
            ]

        if not os.path.exists(pptx_file_path):
            print(f"[ERROR] PPTX file not found: {pptx_file_path}")
            return {"success": False, "output_file": pptx_file_path, "image_paths": image_paths, "error": f"PPTX file not found: {pptx_file_path}"}

        success = insert_images_with_titles(
            pptx_file_path=pptx_file_path,
            image_paths=image_paths,
            slide_number=slide_number,
            top_in=top_in,
            img_w_in=img_w_in,
            side_margin_in=side_margin_in,
            title_height_in=title_height_in,
            title_margin_in=title_margin_in,
            title_font_name=title_font_name,
            title_font_size_pt=title_font_size_pt,
            auto_add_slide_if_missing=auto_add_slide_if_missing
        )

        return {
            "success": success,
            "output_file": pptx_file_path,
            "image_paths": image_paths
        }

    except Exception as e:
        print(f"[ERROR] Failed to run image insertion: {str(e)}")
        return {"success": False, "output_file": pptx_file_path, "image_paths": image_paths, "error": str(e)}

if __name__ == "__main__":
    project_name = "华发四季半岛"
    result = run(
        project_name=project_name,
        slide_number=3,
        top_in=4.5,
        img_w_in=1.8,
        side_margin_in=0.5,
        title_height_in=0.35,
        title_margin_in=0.08
    )
    print("\nAdd room images result:", result)