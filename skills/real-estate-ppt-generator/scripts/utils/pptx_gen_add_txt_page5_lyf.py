#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Insert customer analysis text from a text file into a PowerPoint slide (page 5).
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import os

def run(project_name: str, timestamp: str) -> None:
    """
    Insert customer analysis text from a text file into the 5th slide of a PowerPoint presentation.

    Args:
        project_name (str): Name of the project.
        timestamp (str): Timestamp in YYYYMMDD format.

    Returns:
        None
    """
    # 输入和输出路径统一
    base_dir = f"resources/working_data/{project_name}_{timestamp}/processed_data"
    ppt_path = os.path.join(base_dir, f"{project_name}_gemdale_housing_project_template.pptx")
    txt_path = os.path.join(base_dir, f"{project_name}_客户分析.txt")
    output_path = ppt_path  # 输出路径与输入路径相同

    # 读取 txt 内容
    with open(txt_path, "r", encoding="utf-8") as f:
        content = f.read()

    # 打开 ppt
    prs = Presentation(ppt_path)

    # 确保至少有 5 页
    if len(prs.slides) < 5:
        raise ValueError("PPT 页数不足 5 页！")

    slide = prs.slides[4]  # 索引从 0 开始，第 5 页是 index=4

    # 在第 5 页添加文本框
    left = Inches(1)
    top = Inches(1.5)
    width = Inches(8)
    height = Inches(5)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    # 设置字体样式
    p = text_frame.add_paragraph()
    p.text = content
    p.font.size = Pt(14)

    # 保存输出
    prs.save(output_path)
    print(f"✅ 已生成文件: {output_path}")

if __name__ == "__main__":
    # 示例调用
    run("华发四季半岛", "20250911")