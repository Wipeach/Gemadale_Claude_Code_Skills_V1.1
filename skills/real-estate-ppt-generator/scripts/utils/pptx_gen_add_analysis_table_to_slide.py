#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Script to add the analysis results from Excel file as a table to PowerPoint presentation.
"""

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from pathlib import Path
from datetime import datetime
from typing import Dict, Any

def load_analysis_data_from_excel(excel_file_path: str) -> list:
    """
    Load analysis data from the first sheet of Excel file.
    
    Parameters:
    - excel_file_path: str, path to Excel file
    
    Returns:
    - list: Table data as list of lists
    """
    try:
        # Read the first sheet
        df = pd.read_excel(excel_file_path, sheet_name=0)
        
        # Convert DataFrame to list of lists
        table_data = [df.columns.tolist()]  # Header row
        table_data.extend(df.values.tolist())  # Data rows
        
        return table_data
        
    except Exception as e:
        print(f"[ERROR] Failed to load Excel data: {str(e)}")
        return []

def add_analysis_table_to_slide(pptx_file_path: str, excel_file_path: str,
                               slide_number: int = 1, left_position: float = 2.0,
                               top_position: float = 4.5, table_width: float = 5.0,
                               table_height: float = 2.0, font_name: str = "Arial",
                               font_size: int = 10) -> bool:
    """
    Add analysis results from Excel as table to PowerPoint presentation.
    
    Parameters:
    - pptx_file_path: str, path to PowerPoint file
    - excel_file_path: str, path to Excel analysis results
    - slide_number: int, slide number (1-based indexing)
    - left_position: float, left position in inches from left edge
    - top_position: float, top position in inches from top edge
    - table_width: float, table width in inches
    - table_height: float, table height in inches
    - font_name: str, font name for table text (default: "Arial")
    - font_size: int, font size in points (default: 10)
    
    Returns:
    - bool: True if successful, False otherwise
    """
    try:
        # Step 1: Load table data from Excel
        print("[INFO] Loading analysis data from Excel file...")
        table_data = load_analysis_data_from_excel(excel_file_path)
        
        if not table_data:
            print("[ERROR] No table data loaded")
            return False
            
        print(f"[INFO] Loaded {len(table_data)} rows of data")
        
        # Step 2: Load presentation
        if not os.path.exists(pptx_file_path):
            print(f"[ERROR] PowerPoint file not found: {pptx_file_path}")
            return False
            
        prs = Presentation(pptx_file_path)
        
        # Step 3: Get the specified slide
        if slide_number < 1 or slide_number > len(prs.slides):
            print(f"[ERROR] Invalid slide number: {slide_number}")
            return False
            
        slide = prs.slides[slide_number - 1]
        
        # Step 4: Add table to slide
        rows = len(table_data)
        cols = len(table_data[0]) if rows > 0 else 1
        
        table = slide.shapes.add_table(
            rows, cols,
            Inches(left_position),
            Inches(top_position),
            Inches(table_width),
            Inches(table_height)
        ).table
        
        # Step 5: Populate table with data and apply font formatting
        col_width = Inches(table_width / cols)
        row_height = Inches(table_height / rows)
        
        for i, row_data in enumerate(table_data):
            for j, cell_value in enumerate(row_data):
                if j < cols:
                    cell = table.cell(i, j)
                    cell.text = str(cell_value)
                    
                    table.rows[i].height = row_height
                    table.columns[j].width = col_width
                    
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.name = font_name
                        paragraph.font.size = Pt(font_size)
                        if i == 0:
                            paragraph.font.bold = True
                    
                    cell.margin_left = Inches(0.05)
                    cell.margin_right = Inches(0.05)
                    cell.margin_top = Inches(0.05)
                    cell.margin_bottom = Inches(0.05)
        
        # Step 6: Save presentation
        prs.save(pptx_file_path)
        print(f"[SUCCESS] Analysis table added to: {pptx_file_path}")
        print(f"[INFO] Font applied: {font_name}, Size: {font_size}pt")
        
        return True
        
    except Exception as e:
        print(f"[ERROR] Failed to add analysis table: {str(e)}")
        return False

def run(project_name: str, pptx_file_path: str = None, excel_file_path: str = None,
        slide_number: int = 1, left_position: float = 2.0, top_position: float = 4.5,
        table_width: float = 5.0, table_height: float = 2.0,
        font_name: str = "Arial", font_size: int = 10) -> Dict[str, Any]:
    """
    Run the analysis table addition with a given project name and optional parameters.
    
    Parameters:
    - project_name: str, name of the project
    - pptx_file_path: str, path to PowerPoint file (optional, defaults to specified template)
    - excel_file_path: str, path to Excel analysis results (optional, defaults to deal analysis output)
    - slide_number: int, slide number (1-based indexing)
    - left_position: float, left position in inches from left edge
    - top_position: float, top position in inches from top edge
    - table_width: float, table width in inches
    - table_height: float, table height in inches
    - font_name: str, font name for table text
    - font_size: int, font size in points
    
    Returns:
    - dict: Contains success status, output file path, and table data
    """
    try:
        # Set default file paths
        timestamp = datetime.now().strftime("%Y%m%d")
        if pptx_file_path is None:
            pptx_file_path = str(Path(f"resources/working_data/{project_name}_{timestamp}/processed_data/{project_name}_gemdale_housing_project_template.pptx"))
        if excel_file_path is None:
            excel_file_path = str(Path(f"resources/working_data/{project_name}_{timestamp}/processed_data/{project_name}_成交分析结果.xlsx"))
        
        # Verify files exist
        if not os.path.exists(pptx_file_path):
            print(f"[ERROR] PPTX file not found: {pptx_file_path}")
            return {"success": False, "output_file": pptx_file_path, "table_data": [], "error": f"PPTX file not found: {pptx_file_path}"}
        if not os.path.exists(excel_file_path):
            print(f"[ERROR] Excel file not found: {excel_file_path}")
            return {"success": False, "output_file": pptx_file_path, "table_data": [], "error": f"Excel file not found: {excel_file_path}"}
        
        # Add analysis table
        success = add_analysis_table_to_slide(
            pptx_file_path=pptx_file_path,
            excel_file_path=excel_file_path,
            slide_number=slide_number,
            left_position=left_position,
            top_position=top_position,
            table_width=table_width,
            table_height=table_height,
            font_name=font_name,
            font_size=font_size
        )
        
        # Load table data for return
        table_data = load_analysis_data_from_excel(excel_file_path)
        
        return {
            "success": success,
            "output_file": pptx_file_path,
            "table_data": table_data if table_data else []
        }
        
    except Exception as e:
        print(f"[ERROR] Failed to run analysis table addition: {str(e)}")
        return {"success": False, "output_file": pptx_file_path, "table_data": [], "error": str(e)}

if __name__ == "__main__":
    # For testing purposes, use a default project name
    project_name = "华发四季半岛"
    result = run(
        project_name=project_name,
        slide_number=2,
        left_position=2.0,
        top_position=4.5,
        table_width=5.0,
        table_height=2.0
    )
    print("\nAdd analysis table result:", result)