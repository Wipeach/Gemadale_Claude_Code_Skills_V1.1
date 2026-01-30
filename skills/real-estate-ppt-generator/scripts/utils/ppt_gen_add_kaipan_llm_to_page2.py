#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
读取 resources/working_data/{project_name}_{timestamp}/processed_data/{project_name}_基本信息.txt，
使用 SiliconFlow (OpenAI 兼容) 客户端做一句话开盘摘要，并将该一句话插入到 PPT 第2页顶部（居中显示）。
提供 run(project_name, timestamp=None, pptx_path=None) 接口。
"""

from pathlib import Path
from datetime import datetime
import os
import re
import traceback

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor



# --------------------------
# 启发式摘要（当没有可用 LLM 时回退）
# --------------------------
def heuristic_kaipan_summary(text: str) -> str:
    if not text:
        return "未找到开盘信息。"
    # 找包含“开盘”关键词的句子
    sentences = re.split(r'[。；;\n]+', text)
    keywords = ['开盘', '首开', '开盘时间', '开盘价', '加推', '批次', '首推', '开盘均价', '认筹', '开盘现场']
    candidates = []
    for s in sentences:
        if any(k in s for k in keywords):
            s_clean = s.strip()
            if s_clean:
                candidates.append(s_clean)
    if candidates:
        out = candidates[0].strip()
        return out if len(out) <= 200 else out[:200] + '…'
    # 回退：全文前 120 字
    plain = re.sub(r'\s+', ' ', text).strip()
    if not plain:
        return "未找到开盘信息。"
    return (plain[:120] + '…') if len(plain) > 120 else plain

# --------------------------
# 使用 SiliconFlow (OpenAI 兼容) 客户端进行一句话摘要
# --------------------------
def llm_kaipan_summary(text: str) -> str:
    """
    使用 SiliconFlow API（兼容 OpenAI）生成一句话的开盘摘要。
    若调用失败或不可用，抛出异常由调用方回退到启发式。
    """
    # 设置SiliconFlow API客户端
    # 注意：API密钥应通过环境变量或配置文件提供，此处保留占位符或从 env 获取
    try:
        # 尝试导入新版 OpenAI SDK 中的 OpenAI client
        from openai import OpenAI as OpenAIClient
    except Exception:
        # 旧版 openai 包也可能存在，但我们优先尝试上面那种；如果没有则尝试常规 openai
        try:
            import openai as _openai_mod
            OpenAIClient = None
        except Exception:
            OpenAIClient = None

    api_key = os.getenv("SILICONFLOW_API_KEY", "sk-ykloduxdazjstefqarmswtetrtafvvalaxtkqhxeldvsogtt")
    base_url = "https://api.siliconflow.cn/v1"  # SiliconFlow API基础URL

    # 构造 prompt：要求一句话摘要（中文）
    prompt = (
        "请用一句话（中文）总结下面文本中关于“开盘/开盘信息”的关键信息，"
        "尽量包含时间、批次、价格或重要节点。如果文本中没有开盘相关内容，请直接返回“未找到开盘信息”。"
        "输出不要包含多余解释，仅返回一句话。\n\n文本:\n" + text
    )

    # 尝试使用新版 OpenAIClient 接口（如果可用）
    if OpenAIClient is not None:
        try:
            client = OpenAIClient(api_key=api_key, base_url=base_url)
            # 调用 chat completions（OpenAI 新 SDK 的常见调用方式）
            resp = client.chat.completions.create(
                model="deepseek-ai/DeepSeek-R1",
                messages=[{"role": "user", "content": prompt}],
                temperature=0.1,
                max_tokens=150
            )
            # 解析响应兼容不同 SDK 返回格式
            content = None
            try:
                content = resp.choices[0].message.content
            except Exception:
                try:
                    content = resp['choices'][0]['message']['content']
                except Exception:
                    content = str(resp)
            if not content or not str(content).strip():
                raise RuntimeError("LLM 返回为空")
            return str(content).strip()
        except Exception as e:
            # 抛出异常以让调用方回退启发式
            print("[WARN] SiliconFlow (OpenAIClient) 调用失败:", e)
            # 可打印堆栈便于调试
            # traceback.print_exc()
            raise

    # 否则尝试使用旧的 openai 包（如果存在）
    try:
        import openai as openai_legacy
        # 设置 base url 与 api key 以便走 SiliconFlow 的 API（openai 库允许通过 api_base 覆盖）
        try:
            openai_legacy.api_base = base_url
        except Exception:
            pass
        if hasattr(openai_legacy, "api_key"):
            openai_legacy.api_key = api_key
        else:
            os.environ["OPENAI_API_KEY"] = api_key

        # 尝试 ChatCompletion.create（旧接口）
        resp = openai_legacy.ChatCompletion.create(
            model="gpt-3.5-turbo",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.1,
            max_tokens=150
        )
        content = None
        try:
            content = resp.choices[0].message.content
        except Exception:
            try:
                content = resp['choices'][0]['message']['content']
            except Exception:
                content = str(resp)
        if not content or not str(content).strip():
            raise RuntimeError("LLM 返回为空")
        return str(content).strip()
    except Exception as e:
        print("[WARN] 使用旧 openai 客户端调用 SiliconFlow 失败或不可用:", e)
        # traceback.print_exc()
        raise

# --------------------------
# 保证结果是一句话
# --------------------------
def ensure_one_sentence(s: str) -> str:
    if not s:
        return "未找到开盘信息。"
    s = s.strip().replace('\n', ' ').replace('\r', ' ')
    parts = re.split(r'[。！？!?]+', s)
    first = parts[0].strip() if parts and parts[0].strip() else s.strip()
    if len(first) > 200:
        first = first[:200] + '…'
    if not re.search(r'[。！？!?]$', first):
        first = first + '。'
    return first

# --------------------------
# 主接口：run
# --------------------------
def run(project_name: str, timestamp: str = None, pptx_path: str = None) -> dict:
    """
    读取基本信息文本，总结开盘信息一句话并插入到 PPT 第2页顶部。
    参数：
      - project_name: 项目名称（用于构造路径）
      - timestamp: YYYYMMDD，可选，默认今天
      - pptx_path: 可选，指定要修改的 PPTX 路径；若 None 则使用默认路径 resources/working_data/{project}_{timestamp}/processed_data/{project}_gemdale_housing_project_template.pptx
    返回：
      dict 包含 success(bool)、summary(str)、pptx_path(str) 或 error 信息
    """
    try:
        if not timestamp:
            timestamp = datetime.now().strftime("%Y%m%d")
        base = Path("resources") / "working_data" / f"{project_name}_{timestamp}"
        info_file = base / f"{project_name}_基本信息.txt"
        if pptx_path is None:
            pptx_path = base / f"processed_data/{project_name}_gemdale_housing_project_template.pptx"
        else:
            pptx_path = Path(pptx_path)

        if not info_file.exists():
            return {"success": False, "error": f"基本信息文件不存在: {info_file}"}
        raw_text = info_file.read_text(encoding="utf-8", errors="ignore")
        if not raw_text.strip():
            return {"success": False, "error": f"基本信息文件为空: {info_file}"}

        # 优先尝试 SiliconFlow LLM（如果可用），否则回退到启发式
        summary_raw = None
        try:
            summary_raw = llm_kaipan_summary(raw_text)
        except Exception as e:
            print("[INFO] LLM 摘要不可用或失败，使用启发式摘要。详情：", e)
            summary_raw = heuristic_kaipan_summary(raw_text)

        # 确保是一句话的形式
        summary = ensure_one_sentence(summary_raw)

        # 打开 PPT 并插入文本
        if not pptx_path.exists():
            return {"success": False, "error": f"PPTX 文件不存在: {pptx_path}"}

        prs = Presentation(str(pptx_path))

        # 确保至少有 2 页
        while len(prs.slides) < 2:
            prs.slides.add_slide(prs.slide_layouts[6])

        slide = prs.slides[1]  # 第二页（索引 1）

        # 在第2页顶部插入文本框（居中）
        slide_width = prs.slide_width
        left = Inches(0.6)
        top = Inches(0.75)
        width = slide_width - Inches(1.2)
        height = Inches(0.9)

        # 直接添加一个新的文本框以避免覆盖现有元素
        tx = slide.shapes.add_textbox(left, top, width, height)
        tf = tx.text_frame
        tf.word_wrap = True
        tf.text = ""  # 清空默认内容

        p = tf.paragraphs[0]
        p.text = summary
        p.alignment = PP_ALIGN.CENTER
        # 字体样式（按需可改）
        p.font.name = '黑体'
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 51, 102)

        # 保存 PPT（覆盖原文件）
        prs.save(str(pptx_path))

        return {"success": True, "summary": summary, "pptx_path": str(pptx_path)}
    except Exception as e:
        traceback.print_exc()
        return {"success": False, "error": str(e)}

# 如果直接运行，会以 "华发四季半岛" 做演示（请在真实 pipeline 中调用 run(project_name)）
if __name__ == "__main__":
    res = run("华发四季半岛")
    print(res)
