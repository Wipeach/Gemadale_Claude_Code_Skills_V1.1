# utils/pptx_gen_add_pie_picture_to_page5.py
#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
将地域来源饼图插入到指定PPT的第5页右侧（在原PPT上直接修改并覆盖保存）
提供函数：run(project_name: str, timestamp: str | None) -> dict
"""

from pathlib import Path
from datetime import datetime
from typing import Dict, Any
from pptx import Presentation
from pptx.util import Inches
import os

# EMU per inch，用于将 pptx 单位转换回英寸显示（可选）
EMU_PER_INCH = 914400.0

def run(project_name: str, timestamp: str = None) -> Dict[str, Any]:
    try:
        if not timestamp:
            timestamp = datetime.now().strftime("%Y%m%d")

        base = Path("resources") / "working_data" / f"{project_name}_{timestamp}" / "processed_data"
        pptx_path = base / f"{project_name}_gemdale_housing_project_template.pptx"
        image_path = base / f"{project_name}_地域来源饼图.png"

        if not pptx_path.exists():
            return {"success": False, "error": f"PPTX 文件不存在: {pptx_path}", "pptx_path": str(pptx_path)}

        if not image_path.exists():
            return {"success": False, "error": f"饼图图片不存在: {image_path}", "image_path": str(image_path)}

        # 打开 PPT（修改原文件）
        prs = Presentation(str(pptx_path))

        # 确保至少 5 页（index 4）
        while len(prs.slides) < 5:
            prs.slides.add_slide(prs.slide_layouts[6])

        slide = prs.slides[4]  # 第5页

        # 计算尺寸（英寸）
        slide_width_in = prs.slide_width / EMU_PER_INCH
        slide_height_in = prs.slide_height / EMU_PER_INCH

        target_width_in = slide_width_in * 0.45
        right_margin_in = 0.4
        left_in = slide_width_in - target_width_in - right_margin_in
        if left_in < 0.2:
            left_in = 0.2
        top_in = 1.0
        available_height_in = slide_height_in - top_in - 0.6

        # 尝试使用 Pillow 保持图片长宽比
        width_in = target_width_in
        height_in = None
        try:
            from PIL import Image
            with Image.open(str(image_path)) as im:
                img_w, img_h = im.size
                ratio = (img_h / img_w) if img_w != 0 else 1.0
                est_height_in = width_in * ratio
                if est_height_in > available_height_in:
                    height_in = available_height_in
                    width_in = height_in / ratio if ratio != 0 else width_in
                else:
                    height_in = est_height_in
        except Exception:
            # Pillow 不可用或读取失败：让 pptx 根据宽度自动缩放
            height_in = None

        # 插入图片
        if height_in is None:
            pic = slide.shapes.add_picture(str(image_path), Inches(left_in), Inches(top_in), width=Inches(width_in))
            actual_width_in = pic.width / EMU_PER_INCH
            actual_height_in = pic.height / EMU_PER_INCH
        else:
            pic = slide.shapes.add_picture(str(image_path), Inches(left_in), Inches(top_in), width=Inches(width_in), height=Inches(height_in))
            actual_width_in = width_in
            actual_height_in = height_in

        # 覆盖保存
        prs.save(str(pptx_path))

        placement = {
            "left_in": round(left_in, 3),
            "top_in": round(top_in, 3),
            "width_in": round(actual_width_in, 3),
            "height_in": round(actual_height_in, 3),
            "slide_index": 5
        }

        return {
            "success": True,
            "message": "已将饼图插入到第5页右侧并覆盖保存 PPTX。",
            "pptx_path": str(pptx_path),
            "image_path": str(image_path),
            "placement": placement
        }

    except Exception as e:
        return {"success": False, "error": str(e)}
