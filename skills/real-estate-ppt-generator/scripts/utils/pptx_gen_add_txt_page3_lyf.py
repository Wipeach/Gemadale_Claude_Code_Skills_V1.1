#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Insert text and images from a text file into a PPT slide (page 3).
Parses household analysis text and image paths, adds text to a textbox and images below.
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, List

def parse_txt_file(file_path: str) -> List[Dict[str, str]]:
    """
    Parse text file to extract household analysis data and image paths.
    
    Parameters:
    - file_path: str, path to text file
    
    Returns:
    - List[Dict[str, str]]: List of dictionaries with image_path and overall_evaluation
    """
    try:
        analysis_data = []
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
            sections = content.split("="*50 + "\n\n")
            
            for section in sections:
                if not section.strip():
                    continue
                path_match = re.search(r"户型图路径:\s*(.+)", section)
                image_path = path_match.group(1).strip() if path_match else "未知路径"
                eval_match = re.search(r"###\s*总体评价\s*\n(.+?)(?=\n===|\Z)", section, re.DOTALL)
                overall_evaluation = eval_match.group(1).strip() if eval_match else "（未找到总体评价）"
                
                analysis_data.append({
                    "image_path": image_path,
                    "overall_evaluation": overall_evaluation
                })
        
        return analysis_data
    except Exception as e:
        print(f"[ERROR] Failed to parse text file {file_path}: {str(e)}")
        return []

def insert_text_and_images(
    pptx_file_path: str,
    txt_file_path: str,
    slide_number: int = 3,
    text_left_in: float = 0.0,
    text_top_in: float = 0.5,
    text_width_in: float = 10.0,
    text_height_in: float = 3.5,
    img_top_in: float = 4.5,
    img_width_in: float = 1.5,
    img_height_in: float = 1.0,
    total_width_in: float = 10.0,
    text_font_size_pt: int = 14,
) -> bool:
    """
    Insert text and images into a PPT slide.
    
    Parameters:
    - pptx_file_path: str, path to PowerPoint file
    - txt_file_path: str, path to text file with household analysis
    - slide_number: int, slide number (1-based indexing)
    - text_left_in: float, left position of textbox in inches
    - text_top_in: float, top position of textbox in inches
    - text_width_in: float, width of textbox in inches
    - text_height_in: float, height of textbox in inches
    - img_top_in: float, top position of images in inches
    - img_width_in: float, width of each image in inches
    - img_height_in: float, height of each image in inches
    - total_width_in: float, total width for image layout in inches
    - text_font_size_pt: int, font size for text in points
    
    Returns:
    - bool: True if successful, False otherwise
    """
    try:
        if not os.path.exists(pptx_file_path):
            print(f"[ERROR] PPTX file not found: {pptx_file_path}")
            return False
        if not os.path.exists(txt_file_path):
            print(f"[ERROR] Text file not found: {txt_file_path}")
            return False

        prs = Presentation(pptx_file_path)

        if slide_number < 1:
            print(f"[ERROR] Invalid slide number: {slide_number}")
            return False
        while slide_number > len(prs.slides):
            blank_layout = prs.slide_layouts[6]
            prs.slides.add_slide(blank_layout)

        if slide_number > len(prs.slides):
            print(f"[ERROR] Invalid slide number (exceeds slide count): {slide_number}")
            return False

        slide = prs.slides[slide_number - 1]

        analysis_data = parse_txt_file(txt_file_path)
        if not analysis_data:
            print("[WARNING] No analysis data parsed, saving PPT without changes.")
            prs.save(pptx_file_path)
            print(f"[INFO] Saved PPT without text/images to: {pptx_file_path}")
            return True

        textbox = slide.shapes.add_textbox(
            Inches(text_left_in),
            Inches(text_top_in),
            Inches(text_width_in),
            Inches(text_height_in)
        )
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        for idx, data in enumerate(analysis_data, 1):
            p = text_frame.add_paragraph()
            p.text = f"户型 {idx}: {data['overall_evaluation']}"
            p.font.size = Pt(text_font_size_pt)
            p.font.name = "Arial"
            p.space_after = Pt(10)
            p.level = 0
            p.alignment = PP_ALIGN.LEFT

        img_count = len(analysis_data)
        if img_count > 0:
            img_spacing = total_width_in / img_count
            for idx, data in enumerate(analysis_data):
                img_path = data["image_path"]
                if os.path.exists(img_path):
                    img_left = Inches(idx * img_spacing + (img_spacing - img_width_in) / 2)
                    slide.shapes.add_picture(
                        img_path,
                        img_left,
                        Inches(img_top_in),
                        Inches(img_width_in),
                        Inches(img_height_in)
                    )
                    print(f"[INFO] Inserted image: {img_path} at left={img_left/Inches(1):.2f}in")
                else:
                    print(f"[WARNING] Image not found: {img_path}")

        prs.save(pptx_file_path)
        print(f"[SUCCESS] Saved PPT with text and images to: {pptx_file_path}")
        return True

    except Exception as e:
        print(f"[ERROR] Failed to insert text and images: {str(e)}")
        return False

def run(
    project_name: str,
    pptx_file_path: str = None,
    txt_file_path: str = None,
    slide_number: int = 3,
    text_left_in: float = 0.0,
    text_top_in: float = 0.5,
    text_width_in: float = 10.0,
    text_height_in: float = 3.5,
    img_top_in: float = 4.5,
    img_width_in: float = 1.5,
    img_height_in: float = 1.0,
    total_width_in: float = 10.0,
    text_font_size_pt: int = 14
) -> Dict[str, Any]:
    """
    Run the text and image insertion on page 3 for a given project name.
    
    Parameters:
    - project_name: str, name of the project
    - pptx_file_path: str, path to PowerPoint file (optional, defaults to specified template)
    - txt_file_path: str, path to text file (optional, defaults to household analysis output)
    - slide_number: int, slide number (1-based indexing)
    - text_left_in: float, left position of textbox in inches
    - text_top_in: float, top position of textbox in inches
    - text_width_in: float, width of textbox in inches
    - text_height_in: float, height of textbox in inches
    - img_top_in: float, top position of images in inches
    - img_width_in: float, width of each image in inches
    - img_height_in: float, height of each image in inches
    - total_width_in: float, total width for image layout in inches
    - text_font_size_pt: int, font size for text in points
    
    Returns:
    - dict: Contains success status, output file path, and analysis data
    """
    try:
        timestamp = datetime.now().strftime("%Y%m%d")
        if pptx_file_path is None:
            pptx_file_path = str(Path(f"resources/working_data/{project_name}_{timestamp}/processed_data/{project_name}_gemdale_housing_project_template.pptx"))
        if txt_file_path is None:
            txt_file_path = str(Path(f"resources/working_data/{project_name}_{timestamp}/processed_data/{project_name}_户型分析.txt"))

        if not os.path.exists(pptx_file_path):
            print(f"[ERROR] PPTX file not found: {pptx_file_path}")
            return {"success": False, "output_file": pptx_file_path, "analysis_data": [], "error": f"PPTX file not found: {pptx_file_path}"}
        if not os.path.exists(txt_file_path):
            print(f"[ERROR] Text file not found: {txt_file_path}")
            return {"success": False, "output_file": pptx_file_path, "analysis_data": [], "error": f"Text file not found: {txt_file_path}"}

        analysis_data = parse_txt_file(txt_file_path)
        success = insert_text_and_images(
            pptx_file_path=pptx_file_path,
            txt_file_path=txt_file_path,
            slide_number=slide_number,
            text_left_in=text_left_in,
            text_top_in=text_top_in,
            text_width_in=text_width_in,
            text_height_in=text_height_in,
            img_top_in=img_top_in,
            img_width_in=img_width_in,
            img_height_in=img_height_in,
            total_width_in=total_width_in,
            text_font_size_pt=text_font_size_pt
        )

        return {
            "success": success,
            "output_file": pptx_file_path,
            "analysis_data": analysis_data
        }

    except Exception as e:
        print(f"[ERROR] Failed to run text and image insertion: {str(e)}")
        return {"success": False, "output_file": pptx_file_path, "analysis_data": [], "error": str(e)}

if __name__ == "__main__":
    project_name = "华发四季半岛"
    result = run(
        project_name=project_name,
        slide_number=3,
        text_left_in=0.0,
        text_top_in=0.5,
        text_width_in=10.0,
        text_height_in=3.5,
        img_top_in=4.5,
        img_width_in=1.5,
        img_height_in=1.0
    )
    print("\nAdd text and images result:", result)