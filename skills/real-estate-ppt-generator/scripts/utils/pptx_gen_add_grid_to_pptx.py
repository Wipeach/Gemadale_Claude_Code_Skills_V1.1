from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import sys
import os

def add_grid_to_slide(presentation_path, output_path=None):
    """
    Add a 1-inch spaced grid to the first slide of a PowerPoint presentation
    """
    try:
        # Load the presentation
        prs = Presentation(presentation_path)
        
        if len(prs.slides) == 0:
            print("Error: Presentation has no slides")
            return False
            
        # Get the first slide
        slide = prs.slides[0]
        
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Convert dimensions to inches (1 inch = 914400 EMUs)
        width_inches = slide_width / 914400
        height_inches = slide_height / 914400
        
        print(f"Slide dimensions: {width_inches:.2f} x {height_inches:.2f} inches")
        
        # Main grid line color (light gray)
        grid_color = RGBColor(200, 200, 200)
        # Finer grid color (light blue)
        fine_grid_color = RGBColor(173, 216, 230)  # Light blue
        
        # Add main vertical lines (1-inch spacing)
        for x in range(0, int(width_inches) + 1):
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(x),
                Inches(0),
                Inches(0.01),  # Thin line
                slide_height
            )
            line.fill.solid()
            line.fill.fore_color.rgb = grid_color
            line.line.fill.background()
        
        # Add main horizontal lines (1-inch spacing)
        for y in range(0, int(height_inches) + 1):
            line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                Inches(y),
                slide_width,
                Inches(0.01)  # Thin line
            )
            line.fill.solid()
            line.fill.fore_color.rgb = grid_color
            line.line.fill.background()
            
        # Add finer vertical lines (1/3 inch spacing)
        for x in range(0, int(width_inches * 3) + 1):
            x_pos = x / 3.0
            # Skip positions where main grid lines exist
            if x % 3 != 0 and x_pos <= width_inches:
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(x_pos),
                    Inches(0),
                    Inches(0.005),  # Thinner line
                    slide_height
                )
                line.fill.solid()
                line.fill.fore_color.rgb = fine_grid_color
                line.line.fill.background()
        
        # Add finer horizontal lines (1/3 inch spacing)
        for y in range(0, int(height_inches * 3) + 1):
            y_pos = y / 3.0
            # Skip positions where main grid lines exist
            if y % 3 != 0 and y_pos <= height_inches:
                line = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(0),
                    Inches(y_pos),
                    slide_width,
                    Inches(0.005)  # Thinner line
                )
                line.fill.solid()
                line.fill.fore_color.rgb = fine_grid_color
                line.line.fill.background()
        
        # Add position labels at intersections
        for x in range(0, int(width_inches) + 1):
            for y in range(0, int(height_inches) + 1):
                # Add text box for coordinates
                textbox = slide.shapes.add_textbox(
                    Inches(x) - Inches(0.2),  # Position slightly left of intersection
                    Inches(y) - Inches(0.1),  # Position slightly above intersection
                    Inches(0.4),  # Width for the text
                    Inches(0.2)   # Height for the text
                )
                
                text_frame = textbox.text_frame
                text_frame.text = f"({x},{y})"
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                text_frame.paragraphs[0].font.size = Pt(6)
                text_frame.paragraphs[0].font.color.rgb = RGBColor(100, 100, 100)
                text_frame.margin_left = Pt(0)
                text_frame.margin_right = Pt(0)
                text_frame.margin_top = Pt(0)
                text_frame.margin_bottom = Pt(0)
        
        # Determine output path
        if output_path is None:
            base_path = os.path.splitext(presentation_path)[0]
            output_path = f"{base_path}_with_grid.pptx"
        
        # Save the presentation
        prs.save(output_path)
        print(f"Grid added successfully!")
        print(f"Output file: {output_path}")
        return True

    except Exception as e:
        print(f"Error adding grid: {e}")
        import traceback
        traceback.print_exc()
        return False

def main():
    """
    Main execution function
    """
    print("Adding 1-inch grid to PowerPoint presentation...")
    print("=" * 50)
    
    # Check if input file exists
    input_file = "gemdale_slide_master_template_with_analysis.pptx"
    
    if not os.path.exists(input_file):
        print(f"Error: Input file '{input_file}' not found!")
        print("Please run 'create_gemdale_slide_master.py' first to create the template file.")
        return
    
    # Add grid to the presentation
    success = add_grid_to_slide(input_file)
    
    if success:
        print("=" * 50)
        print("Grid features added:")
        print("- Horizontal and vertical lines at 1-inch intervals (main grid)")
        print("- Finer grid with 1/3 inch spacing between main lines")
        print("- Position labels (x,y) at main grid intersections only")
        print("- Light gray color for main grid, light blue for finer grid")
        print("- Thinner lines for finer grid elements")
    else:
        print("Failed to add grid. Please check the error messages above.")

if __name__ == "__main__":
    main()