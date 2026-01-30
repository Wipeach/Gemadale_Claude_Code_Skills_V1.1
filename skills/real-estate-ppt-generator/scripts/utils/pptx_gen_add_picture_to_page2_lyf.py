#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Insert an image into a PPT slide at a given position (in inches),
and add a centered title above the image using project_name.
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image
from datetime import datetime
from pathlib import Path
from typing import Dict, Any, Optional, Tuple, List

EMU_PER_INCH = 914400


def _compute_left_top_by_anchor(left_in: float, top_in: float, w_in: float, h_in: float, anchor: str) -> Tuple[float, float]:
    """
    Convert an anchor point with given (left_in, top_in) into the
    actual top-left coordinates for the picture.
    anchor in {"top_left","top_right","center","bottom_left","bottom_right"}
    """
    anchor = (anchor or "top_left").lower()
    if anchor == "top_left":
        return left_in, top_in
    elif anchor == "top_right":
        return left_in - w_in, top_in
    elif anchor == "center":
        return left_in - w_in / 2.0, top_in - h_in / 2.0
    elif anchor == "bottom_left":
        return left_in, top_in - h_in
    elif anchor == "bottom_right":
        return left_in - w_in, top_in - h_in
    else:
        return left_in, top_in


def insert_image_with_title(
    pptx_file_path: str,
    image_path: str,
    project_name: str,
    slide_number: int = 1,
    left_in: float = 6.5,
    top_in: float = 5.8,
    anchor: str = "top_left",
    image_width_in: float = 5.5,
    title_font_name: str = "Arial",
    title_font_size_pt: int = 16,
    title_height_in: float = 0.45,
    title_margin_in: float = 0.08,
    auto_add_slide_if_missing: bool = True,
    output_file: Optional[str] = None
) -> str:
    """
    Insert image and title into a PPT slide at absolute position.

    Returns:
      pptx_file_path (string) if success, otherwise empty string.
    """
    try:
        if not os.path.exists(pptx_file_path):
            print(f"[ERROR] PPTX file not found: {pptx_file_path}")
            return ""

        if not os.path.exists(image_path):
            print(f"[ERROR] Image file not found: {image_path}")
            return ""

        prs = Presentation(pptx_file_path)

        if slide_number < 1:
            print(f"[ERROR] Invalid slide number: {slide_number}")
            return ""

        # add blank slides if needed
        while auto_add_slide_if_missing and slide_number > len(prs.slides):
            blank_layout = prs.slide_layouts[6]
            prs.slides.add_slide(blank_layout)

        if slide_number > len(prs.slides):
            print(f"[ERROR] Invalid slide number (exceeds slide count): {slide_number}")
            return ""

        slide = prs.slides[slide_number - 1]

        with Image.open(image_path) as img:
            img_w_px, img_h_px = img.size
        if img_w_px <= 0 or img_h_px <= 0:
            print("[ERROR] Image size invalid.")
            return ""

        aspect = img_h_px / float(img_w_px)
        target_w_in = float(image_width_in)
        target_h_in = target_w_in * aspect

        tl_left_in, tl_top_in = _compute_left_top_by_anchor(left_in, top_in, target_w_in, target_h_in, anchor)

        pic = slide.shapes.add_picture(
            image_path,
            Inches(tl_left_in),
            Inches(tl_top_in),
            width=Inches(target_w_in),
            height=Inches(target_h_in)
        )
        # ensure position/size is set explicitly
        pic.left = Inches(tl_left_in)
        pic.top = Inches(tl_top_in)
        pic.width = Inches(target_w_in)
        pic.height = Inches(target_h_in)

        # title textbox above picture
        title_height_emu = Inches(title_height_in)
        title_margin_emu = Inches(title_margin_in)
        title_top_emu = pic.top - title_height_emu - title_margin_emu
        if title_top_emu < Inches(0.1):
            title_top_emu = Inches(0.1)

        textbox = slide.shapes.add_textbox(
            pic.left,
            title_top_emu,
            pic.width,
            title_height_emu
        )
        tf = textbox.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = f"{project_name}分户型销售情况"
        p.alignment = PP_ALIGN.CENTER
        # safe access to run
        if len(p.runs) == 0:
            # ensure at least one run exists
            r = p.add_run()
            r.text = p.text
        run_obj = p.runs[0]
        run_obj.font.name = title_font_name
        run_obj.font.size = Pt(title_font_size_pt)
        run_obj.font.bold = True

        # Use pptx_file_path as output path
        output_file = pptx_file_path

        # Ensure parent exists
        os.makedirs(os.path.dirname(output_file), exist_ok=True)

        prs.save(output_file)

        slide_w_in = prs.slide_width / EMU_PER_INCH
        slide_h_in = prs.slide_height / EMU_PER_INCH
        print(f"[INFO] Slide size: {slide_w_in:.2f}in x {slide_h_in:.2f}in")
        print(f"[INFO] Image anchor='{anchor}' at ({left_in:.2f}, {top_in:.2f}) in")
        print(f"[INFO] Picture top-left -> ({tl_left_in:.2f}, {tl_top_in:.2f}) in ; size -> {target_w_in:.2f} x {target_h_in:.2f} in")
        print(f"[SUCCESS] Saved PPT with image and title to: {output_file}")
        return output_file

    except Exception as e:
        print(f"[ERROR] Failed to insert image and title: {e}")
        return ""


def run(
    project_name: str,
    pptx_file_path: Optional[str] = None,
    image_path: Optional[str] = None,
    slide_number: int = 2,
    left_in: float = 8.5,
    top_in: float = 2.8,
    anchor: str = "top_left",
    image_width_in: float = 4.5,
    title_font_name: str = "Arial",
    title_font_size_pt: int = 12,
    title_height_in: float = 0.45,
    title_margin_in: float = 0.08,
    auto_add_slide_if_missing: bool = True,
    timestamp: Optional[str] = None
) -> Dict[str, Any]:
    """
    Run the image and title insertion.

    Returns:
      dict with keys: success (bool), output_file (str), image_path (str), input_pptx (str)
    """
    try:
        if timestamp is None:
            timestamp = datetime.now().strftime("%Y%m%d")

        base_dir = Path(f"resources/working_data/{project_name}_{timestamp}/processed_data")
        # if pptx_file_path not provided, use specified template
        if pptx_file_path is None:
            pptx_file_path = str(base_dir / f"{project_name}_gemdale_housing_project_template.pptx")

        # default image path if not provided
        if image_path is None:
            image_path = str(base_dir / f"{project_name}_成交结果分析混合图与表.png")

        if not os.path.exists(pptx_file_path):
            return {"success": False, "output_file": "", "image_path": image_path, "input_pptx": pptx_file_path, "error": f"PPTX file not found: {pptx_file_path}"}
        if not os.path.exists(image_path):
            return {"success": False, "output_file": "", "image_path": image_path, "input_pptx": pptx_file_path, "error": f"Image file not found: {image_path}"}

        # Use same path for output
        output_file = pptx_file_path

        # Insert image and title; get output path (string) back
        result_path = insert_image_with_title(
            pptx_file_path=pptx_file_path,
            image_path=image_path,
            project_name=project_name,
            slide_number=slide_number,
            left_in=left_in,
            top_in=top_in,
            anchor=anchor,
            image_width_in=image_width_in,
            title_font_name=title_font_name,
            title_font_size_pt=title_font_size_pt,
            title_height_in=title_height_in,
            title_margin_in=title_margin_in,
            auto_add_slide_if_missing=auto_add_slide_if_missing,
            output_file=output_file
        )

        if not result_path:
            return {"success": False, "output_file": "", "image_path": image_path, "input_pptx": pptx_file_path, "error": "Failed to insert image and save output."}

        return {"success": True, "output_file": result_path, "image_path": image_path, "input_pptx": pptx_file_path}

    except Exception as e:
        return {"success": False, "output_file": "", "image_path": image_path or "", "input_pptx": pptx_file_path or "", "error": str(e)}


if __name__ == "__main__":
    project_name = "华发四季半岛"
    res = run(project_name=project_name)
    print("Result:", res)