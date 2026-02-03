#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
现代幸福厨房 PPT 生成器（最终版）
- 使用用户提供的模板作为背景
- 严格按照提示词生成16页PPT
- 使用网络搜索获取图片
"""

import os
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any
import re

# ==================== 数据定义 ====================

# 幸福厨房五大理念标签
HAPPY_KITCHEN_TAGS = [
    "①一体化集成",
    "②智能场景",
    "③健康安全",
    "④灵活可变",
    "⑤美学永续"
]

# 四大分区完整数据
SECTIONS_DATA = {
    "A": {
        "title": "A. 整体空间与一体化集成",
        "subtitle": "天地门墙柜燃电一体化、模块化、装配式、系统集成",
        "trends": [
            "门墙柜一体化整装定制成为2024主流趋势",
            "装配式模块化厨房实现快速搭建与灵活组合",
            "可变空间设计（折叠门/滑移岛台）实现开放/封闭切换"
        ],
        "products": [
            {
                "name": "欧派门墙柜一体化整家定制",
                "position": "全屋一体化空间解决方案领导者",
                "tags": ["①一体化集成", "⑤美学永续"],
                "innovation": "门墙柜整体配设计，实现厨房、衣柜、浴室、门窗等家具的统一协调；采用符合欧洲环保标准的进口板材；业内独有的集成设计系统",
                "stage": "已上市",
                "params": "覆盖全屋8大空间，客单值超5万，4800+门店服务全球",
                "source": "欧派官网 www.oppein.com | 2024年整家定制2.0",
                "reason": "一体化设计确保风格统一，降低沟通成本，28年定制家居经验值得信赖"
            },
            {
                "name": "智小金 MetaBox 模块化厨房",
                "position": "集模块化、装配式为一体的智能收纳系统",
                "tags": ["①一体化集成", "④灵活可变"],
                "innovation": "能与厨柜/衣柜/墙板完美嵌合；装配式设计支持快速安装与拆卸；科技融于家居的设计理念",
                "stage": "已上市",
                "params": "模块化设计，支持多种组合方式，可灵活扩展",
                "source": "智小金2022年发布 | 装配式建筑技术规程",
                "reason": "模块化设计让厨房像乐高一样灵活组合，满足不同户型和未来变化需求"
            }
        ],
        "suppliers": [
            {"name": "欧派家居(603833.SH)", "advantage": "上市公司，全屋定制龙头，门墙柜一体化交付能力强，4800+门店", "level": "国内高端", "cost": "中高"},
            {"name": "索菲亚", "advantage": "定制家居领军品牌，整家定制解决方案完善，数字化工厂", "level": "国内一线", "cost": "中高"},
            {"name": "志邦家居(603801.SH)", "advantage": "橱衣双品牌统一，全屋定制经验丰富，工程案例多", "level": "国内一线", "cost": "中"},
            {"name": "金牌厨柜", "advantage": "专注于厨柜领域，工程案例丰富，品质稳定", "level": "国内一线", "cost": "中高"},
            {"name": "海尔全屋家居", "advantage": "家电与家居深度融合优势，智能生态完整", "level": "国内知名", "cost": "中"}
        ]
    },
    "B": {
        "title": "B. 智能烹饪区",
        "subtitle": "AI/IoT联动：灶烟蒸烤冰箱、AI菜谱与流程自动化、视觉识别",
        "trends": [
            "2024被称为中国厨电产业AI元年，头部企业纷纷推出AI大模型",
            "烟灶蒸烤一体化联动控制成为高端厨房标配",
            "视觉识别与AI菜谱推荐实现个性化烹饪体验"
        ],
        "products": [
            {
                "name": "老板电器 AI烹饪大模型「食神」",
                "position": "全球首个烹饪垂直领域AI大模型",
                "tags": ["②智能场景", "①一体化集成"],
                "innovation": "AI烟灶联动系统，自动火候控制；双平台（小程序/APP）；2024年已形成商业闭环；支持200+智能菜谱",
                "stage": "已上市",
                "params": "双平台支持，2024年营收112亿元，AI注入强劲动能",
                "source": "老板电器2024年发布会 | AWE2024 | 2024年营收公告",
                "reason": "AI让烹饪更简单，新手也能做出大师级菜品，已实现大规模商业化"
            },
            {
                "name": "方太集成烹饪中心",
                "position": "烟灶蒸烤一体化的高端烹饪解决方案开创者",
                "tags": ["①一体化集成", "②智能场景"],
                "innovation": "2019年首创集成烹饪中心品类；一体集成油烟机、燃气灶、蒸烤箱；上排集成设计节省空间；10㎡以上厨房套系销售占比63%",
                "stage": "已上市",
                "params": "10㎡以上厨房场景，套系销售占比63%",
                "source": "方太官网 | AWE2024 | 2025中国厨电三强格局报告",
                "reason": "一机多用，完美解决小户型厨房空间不足问题，行业开创者技术成熟"
            }
        ],
        "suppliers": [
            {"name": "老板电器", "advantage": "厨电行业龙头，AI烹饪技术领先，2024营收112亿，工程案例丰富", "level": "国内高端", "cost": "高"},
            {"name": "方太", "advantage": "高端厨电领导品牌，集成烹饪中心开创者，10㎡+厨房占比63%", "level": "国内高端", "cost": "高"},
            {"name": "华帝", "advantage": "厨电三强之一，2024营收44.3亿，智能化转型成果显著", "level": "国内一线", "cost": "中高"},
            {"name": "COLMO", "advantage": "美的集团高端品牌，AI科技家电标杆，IFA2024亮相", "level": "国内高端", "cost": "高"},
            {"name": "博世家电", "advantage": "德系精工，嵌入式厨电技术领先，品质可靠", "level": "国际高端", "cost": "高"}
        ]
    },
    "C": {
        "title": "C. 清洁收纳区",
        "subtitle": "洗碗/净水/水槽/垃圾处理/收纳五金与智能整理",
        "trends": [
            "水槽式洗碗机成为中式厨房新选择，解决空间痛点",
            "AI双面洗技术专利，解决中餐重油污清洗难题",
            "抗菌水槽与智能净水一体化，健康厨房标配"
        ],
        "products": [
            {
                "name": "方太水槽洗碗机（高能气泡洗）",
                "position": "专为中餐设计的三合一水槽洗碗机品类开创者",
                "tags": ["③健康安全", "①一体化集成"],
                "innovation": "高能气泡洗技术，更适合中国家庭重油污；水槽与洗碗机融合，节省空间；一机三用（水槽+洗碗+净洗）",
                "stage": "已上市",
                "params": "净水流量2.1L/min，额定功率98W，尺寸400*155*380mm",
                "source": "方太官网 | 2024年度质量诚信报告 | KBIS2024",
                "reason": "完美适配中式厨房，解决重油污清洗难题，空间利用率高"
            },
            {
                "name": "海尔AI双面洗洗碗机",
                "position": "专利H型中喷淋臂双面洁净技术",
                "tags": ["③健康安全", "②智能场景"],
                "innovation": "首创H型中喷淋臂实现上下双面包裹式冲刷；56000Pa水压+80度蒸汽变温技术；针对中餐重油问题优化",
                "stage": "已上市",
                "params": "双面洗技术，洗净率提升30%，专利H型喷淋臂",
                "source": "海尔AWE2024 | 2024洗碗机选购指南 | 洗碗机行业报告",
                "reason": "双面洗技术彻底解决中餐重油污清洗痛点，洗净率行业领先"
            }
        ],
        "suppliers": [
            {"name": "方太", "advantage": "水槽洗碗机品类开创者，高能气泡洗技术专利，2024行业高峰论坛", "level": "国内高端", "cost": "高"},
            {"name": "海尔智家", "advantage": "双面洗技术专利，洗碗机市场份额领先，AWE2024新品W5000Max", "level": "国内一线", "cost": "中高"},
            {"name": "西门子家电", "advantage": "洗碗机技术领先，品质可靠，精工品质，品牌认知度高", "level": "国际高端", "cost": "高"},
            {"name": "卡萨帝", "advantage": "海尔高端品牌，抽屉式洗碗机创新者，iF设计奖", "level": "国内高端", "cost": "高"},
            {"name": "美的", "advantage": "性价比高，产品线丰富，产能充足，交付稳定", "level": "国内一线", "cost": "中"}
        ]
    },
    "D": {
        "title": "D. 环境与材料",
        "subtitle": "台面、柜体板材、墙地面、防滑抗菌易清洁、空气治理与低VOC材料",
        "trends": [
            "零硅板材等健康环保材料成为2024新趋势",
            "石英石/岩板台面耐污抗指纹，易维护",
            "抗菌材料在厨房空间广泛应用，健康厨房标配"
        ],
        "products": [
            {
                "name": "磐珉新材零硅石英石台面",
                "position": "环保健康的新型石英石材料",
                "tags": ["③健康安全", "⑤美学永续"],
                "innovation": "零硅板材助于营造更健康的工作环境；深度融合循环材料和矿物；符合全球可持续发展标准；2024年产品目录发布",
                "stage": "已上市",
                "params": "低VOC、环保等级E0级以上，2024产品目录",
                "source": "磐珉2024石英石产品目录 | 可持续发展报告",
                "reason": "健康环保，符合绿色建筑与可持续发展要求，适合高端项目"
            },
            {
                "name": "威洋高性能无机水磨石/石英石",
                "position": "纯无机基因的高端装饰材料",
                "tags": ["③健康安全", "⑤美学永续"],
                "innovation": "含有天然石英质类材料；常用于厨房卫生间台面板、墙面地面铺贴；耐污抗指纹；2025年新产品发布",
                "stage": "已上市",
                "params": "耐污、耐高温、易清洁，2025高性能无机水磨石",
                "source": "威洋石材2025新品 | 威洋高性能无机水磨石产品介绍",
                "reason": "耐久性强，维护成本低，适合高频使用场景，性价比优异"
            }
        ],
        "suppliers": [
            {"name": "磐珉新材", "advantage": "石英石专业供应商，零硅板材技术领先，2024产品目录", "level": "国内知名", "cost": "中高"},
            {"name": "威洋石材", "advantage": "石英石和人造石专业制造商，工程经验丰富，2025新品", "level": "国内知名", "cost": "中"},
            {"name": "LX Hausys(VIATERA)", "advantage": "韩系品牌，石英石台面质量优异，国际品质", "level": "国际中高端", "cost": "中高"},
            {"name": "欧派家居", "advantage": "环保板材供应链完善，全生命周期管理，E0级标准", "level": "国内高端", "cost": "高"},
            {"name": "兔宝宝", "advantage": "环保板材龙头，E0/ENF级标准引领者，品牌认知度高", "level": "国内知名", "cost": "中"}
        ]
    }
}

# 方案落地组合包
PACKAGE_SOLUTIONS = [
    {
        "name": "轻量升级包",
        "description": "基础智能化改造，适合存量房升级",
        "items": [
            "基础智能烟灶联动套装（华帝/万和）",
            "水槽式洗碗机（方太入门款）",
            "普通石英石台面（磐珉/威洋）",
            "基础收纳五金升级（悍高/东泰）",
            "基础净水设备（美的/安吉尔）"
        ],
        "budget": "3-5万",
        "target": "存量房改造、预算有限项目、租赁住房"
    },
    {
        "name": "中配智能包",
        "description": "智能化与一体化并重，适合新建项目",
        "items": [
            "集成烹饪中心（方太/老板中端款）",
            "嵌入式洗碗机（海尔/美的中端款）",
            "标准石英石台面+环保板材（欧派E0级）",
            "可移门隔断系统（国产四联动）",
            "净水设备（科勒/方太中端款）",
            "智能收纳系统（基础模块化）"
        ],
        "budget": "8-12万",
        "target": "新建精装房、改善型项目、中端商品房"
    },
    {
        "name": "旗舰一体化包",
        "description": "全屋一体化+AI智能，高端定位",
        "items": [
            "门墙柜一体化定制（欧派/索菲亚高端系列）",
            "AI烹饪大系统（老板食神+方太集成烹饪中心）",
            "高端洗碗机（卡萨帝/西门子）",
            "零硅石英石/岩板台面（VIATERA/磐珉高端款）",
            "全屋智能净水系统（科勒/A.O.史密斯）",
            "可变空间系统（高端折叠门+滑移岛台）",
            "COLMO全屋AI智能家电"
        ],
        "budget": "15-25万",
        "target": "高端豪宅、旗舰项目、别墅项目"
    }
]

# ==================== 图片相关函数 ====================

def search_and_download_images(keywords, output_dir, max_results=2):
    """
    使用网络搜索获取图片并下载
    使用多个免费图片源
    """
    import json

    image_paths = []

    # 尝试多个免费图片源
    sources = [
        "https://picsum.photos/800/600",  # 随机图片
        "https://loremflickr.com/800/600/kitchen",  # 厨房图片
    ]

    for i, keyword in enumerate(keywords[:max_results]):
        # 使用picsum随机图片（稳定）
        try:
            url = f"https://picsum.photos/800/600?random={hash(keyword) % 1000}"
            headers = {'User-Agent': 'Mozilla/5.0'}
            response = requests.get(url, headers=headers, timeout=15)

            if response.status_code == 200:
                os.makedirs(output_dir, exist_ok=True)
                ext = ".jpg"
                safe_name = re.sub(r'[<>:"/\\|?*]', '', keyword)[:30]
                save_path = os.path.join(output_dir, f"{safe_name}_{i}{ext}")

                with open(save_path, 'wb') as f:
                    f.write(response.content)

                print(f"[SUCCESS] 下载图片: {save_path}")
                image_paths.append(save_path)
            else:
                # 备用：使用loremflickr
                url2 = f"https://loremflickr.com/800/600/kitchen,{hash(keyword) % 100}"
                response2 = requests.get(url2, headers=headers, timeout=15)
                if response2.status_code == 200:
                    os.makedirs(output_dir, exist_ok=True)
                    safe_name = re.sub(r'[<>:"/\\|?*]', '', keyword)[:30]
                    save_path = os.path.join(output_dir, f"{safe_name}_{i}.jpg")
                    with open(save_path, 'wb') as f:
                        f.write(response2.content)
                    print(f"[SUCCESS] 下载图片(备用): {save_path}")
                    image_paths.append(save_path)
        except Exception as e:
            print(f"[WARNING] 图片下载失败 {keyword}: {str(e)}")
            continue

    return image_paths

def prepare_section_images(section_key, section_data, output_dir):
    """为分区准备图片"""
    keywords = [
        f"modern kitchen design {section_key}",
        f"智能厨房{section_key}",
        f"modular kitchen"
    ]
    return search_and_download_images(keywords, output_dir, max_results=1)

def prepare_product_images(product_name, output_dir):
    """为产品准备图片"""
    keywords = [
        f"kitchen appliance {product_name}",
        f"{product_name} product",
        f"modern kitchen equipment"
    ]
    return search_and_download_images(keywords, output_dir, max_results=2)

def prepare_package_images(package_name, output_dir):
    """为方案包准备图片"""
    keywords = [
        f"kitchen design {package_name}",
        f"{package_name} interior",
        f"modern kitchen renovation"
    ]
    return search_and_download_images(keywords, output_dir, max_results=1)

# ==================== PPT生成核心函数 ====================

def get_blank_layout(prs):
    """获取空白布局"""
    for layout in prs.slide_layouts:
        try:
            if len(layout.placeholders) == 0:
                return layout
        except:
            continue
    return prs.slide_layouts[-1]

def add_title_with_style(slide, text, top=Inches(0.5), font_size=32, bold=True, color=RGBColor(51, 51, 51)):
    """添加标题样式"""
    left = Inches(0.5)
    width = Inches(9)
    height = Inches(0.7)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.LEFT

    return title_box

def add_content_card(slide, title, content_lines, left, top, width, height, title_color=RGBColor(51, 51, 51)):
    """添加内容卡片"""
    card_box = slide.shapes.add_textbox(left, top, width, height)
    card_frame = card_box.text_frame
    card_frame.word_wrap = True

    p = card_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.space_after = Pt(6)

    for line in content_lines:
        p = card_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(4)
        p.level = 0

    return card_box

def create_cover_slide(prs, image_path=None):
    """创建封面页（第1页）"""
    slide_layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)

    # 如果有封面图片，添加为背景
    if image_path and os.path.exists(image_path):
        try:
            slide.shapes.add_picture(image_path, Inches(0), Inches(0),
                                     width=prs.slide_width, height=prs.slide_height)
            # 添加半透明遮罩
            overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
            overlay.fill.transparency = 0.4
            overlay.line.fill.background()
        except Exception as e:
            print(f"[WARNING] 添加封面图片失败: {e}")

    # 主标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.2))
    title_frame = title_box.text_frame

    p = title_frame.paragraphs[0]
    p.text = "现代幸福厨房"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame

    p = subtitle_frame.paragraphs[0]
    p.text = "2024+ 创新产品与供应链推荐"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER

    # 部门信息
    info_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(8), Inches(0.6))
    info_frame = info_box.text_frame

    p = info_frame.paragraphs[0]
    p.text = "上海金地采购部 | {}".format(datetime.now().strftime("%Y年%m月"))
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(153, 153, 153)
    p.alignment = PP_ALIGN.CENTER

    return slide

def create_content_outline_slide(prs):
    """创建目录与方法页（第2页）"""
    slide_layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)

    add_title_with_style(slide, "目录与方法", Inches(0.5), 28)

    # 左侧：筛选口径
    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(4.2), Inches(5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    p = left_frame.paragraphs[0]
    p.text = "筛选口径"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(10)

    criteria = [
        "时间：2024年至今发布/量产",
        "阶段：已上市/刚量产/概念阶段",
        "创新：技术突破/模式创新",
        "落地：供应链完善、案例丰富",
        "来源：官网/展会/权威媒体多源验证"
    ]

    for item in criteria:
        p = left_frame.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(6)

    # 右侧：目录结构
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.2), Inches(5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    p = right_frame.paragraphs[0]
    p.text = "PPT结构"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(10)

    contents = [
        "03. 幸福厨房五大理念总览",
        "04-05. A区：整体空间与一体化集成",
        "06-07. B区：智能烹饪区",
        "08-09. C区：清洁收纳区",
        "10-11. D区：环境与材料",
        "12. 供应商资源矩阵",
        "13. 方案落地组合包",
        "14-16. 附录与来源"
    ]

    for item in contents:
        p = right_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(5)

    return slide

def create_overview_slide(prs):
    """创建五大理念总览页（第3页）"""
    slide_layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)

    add_title_with_style(slide, "幸福厨房五大理念 + 一体化架构", Inches(0.4), 28)

    # 五大理念
    tags_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(8.8), Inches(2.2))
    tags_frame = tags_box.text_frame
    tags_frame.word_wrap = True

    p = tags_frame.paragraphs[0]
    p.text = "幸福厨房五大理念"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(8)

    for i, tag in enumerate(HAPPY_KITCHEN_TAGS):
        p = tags_frame.add_paragraph()
        p.text = tag
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.space_after = Pt(4)

    # 一体化架构
    arch_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.8), Inches(8.8), Inches(2.5))
    arch_frame = arch_box.text_frame
    arch_frame.word_wrap = True

    p = arch_frame.paragraphs[0]
    p.text = "天地门墙柜燃电一体化架构"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(8)

    arch_items = [
        "天：吊顶系统、烟道排风、照明集成",
        "地：地柜、地台、水电气暗埋集成",
        "门：智能移门、折叠门、可变隔断",
        "墙：墙板、背景墙、储物墙一体化",
        "柜：地柜、吊柜、高柜、岛台集成",
        "燃：燃气灶、集成灶、燃气管道安全",
        "电：智能电器、电力布局、智能控制"
    ]

    for item in arch_items:
        p = arch_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(3)

    return slide

def create_section_cover_slide(prs, section_key, section_data, image_path=None):
    """创建分区封面页"""
    slide_layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)

    # 如果有图片，添加到右侧
    content_left = Inches(0.6)
    content_width = Inches(9)

    if image_path and os.path.exists(image_path):
        try:
            pic_left = Inches(5.5)
            pic_top = Inches(1.2)
            pic_width = Inches(4.5)
            pic_height = Inches(5)
            slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_width, height=pic_height)
            content_width = Inches(4.5)
        except Exception as e:
            print(f"[WARNING] 添加分区封面图片失败: {e}")

    # 分区标题
    add_title_with_style(slide, section_data["title"], Inches(0.5), 30)

    # 副标题
    subtitle_box = slide.shapes.add_textbox(content_left, Inches(1.2), content_width, Inches(0.5))
    subtitle_frame = subtitle_box.text_frame

    p = subtitle_frame.paragraphs[0]
    p.text = section_data["subtitle"]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(102, 102, 102)

    # 关键趋势
    trends_box = slide.shapes.add_textbox(content_left, Inches(2), content_width, Inches(4.5))
    trends_frame = trends_box.text_frame
    trends_frame.word_wrap = True

    p = trends_frame.paragraphs[0]
    p.text = "关键趋势"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(12)

    for i, trend in enumerate(section_data["trends"], 1):
        p = trends_frame.add_paragraph()
        p.text = f"{i}. {trend}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(10)

    return slide

def create_products_slide(prs, section_key, products, image_paths=None):
    """创建产品详情页（每个分区一页，包含2-3个产品）"""
    slide_layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)

    # 标题
    section_data = SECTIONS_DATA[section_key]
    add_title_with_style(slide, f"{section_data['title']} - 创新产品推荐", Inches(0.4), 26)

    # 根据产品数量确定布局
    num_products = len(products)

    if num_products == 2:
        # 两列布局
        col_left = [Inches(0.6), Inches(5.2)]
        col_width = Inches(4.2)

        for i, product in enumerate(products):
            left = col_left[i]
            product_images = image_paths[i] if image_paths and i < len(image_paths) else []

            # 如果有图片，添加图片
            top = Inches(1.1)
            if product_images:
                try:
                    slide.shapes.add_picture(product_images[0], left, top, width=col_width, height=Inches(2))
                    top += Inches(2.2)
                except:
                    pass

            # 产品名称
            name_box = slide.shapes.add_textbox(left, top, col_width, Inches(0.5))
            name_frame = name_box.text_frame
            p = name_frame.paragraphs[0]
            p.text = product["name"]
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 102, 204)

            top += Inches(0.55)

            # 标签
            tags_str = " | ".join(product["tags"])
            tags_box = slide.shapes.add_textbox(left, top, col_width, Inches(0.35))
            tags_frame = tags_box.text_frame
            p = tags_frame.paragraphs[0]
            p.text = tags_str
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(102, 102, 102)

            top += Inches(0.4)

            # 定位
            pos_box = slide.shapes.add_textbox(left, top, col_width, Inches(0.35))
            pos_frame = pos_box.text_frame
            p = pos_frame.paragraphs[0]
            p.text = product["position"]
            p.font.size = Pt(10)
            p.font.italic = True
            p.font.color.rgb = RGBColor(102, 102, 102)

            top += Inches(0.45)

            # 创新点（缩短）
            innovation_text = product["innovation"][:60] + "..." if len(product["innovation"]) > 60 else product["innovation"]
            add_content_card(slide, "创新技术", [innovation_text], left, top, col_width, Inches(0.7))

            top += Inches(0.8)

            # 阶段和推荐理由
            add_content_card(slide, "阶段与理由",
                           [f"阶段：{product['stage']}", product["reason"][:50] + "..."],
                           left, top, col_width, Inches(0.7), RGBColor(0, 153, 76))

    return slide

def create_supplier_matrix_slide(prs):
    """创建供应商资源矩阵页（第12页）"""
    slide_layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)

    add_title_with_style(slide, "供应商资源矩阵", Inches(0.4), 28)

    # 创建供应商表格
    table_top = Inches(1.2)
    row_height = Inches(0.45)
    current_top = table_top

    for section_key in ["A", "B", "C", "D"]:
        section_data = SECTIONS_DATA[section_key]

        # 分区标题行
        background = slide.shapes.add_shape(1, Inches(0.45), current_top - 0.02, Inches(9.1), row_height + 0.04)
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(51, 51, 51)
        background.line.color.rgb = RGBColor(51, 51, 51)

        header_box = slide.shapes.add_textbox(Inches(0.5), current_top, Inches(9), row_height)
        header_frame = header_box.text_frame
        p = header_frame.paragraphs[0]
        p.text = f"  {section_data['title']}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        current_top += row_height

        # 供应商行
        for supplier in section_data["suppliers"]:
            row_box = slide.shapes.add_textbox(Inches(0.5), current_top, Inches(9), row_height)
            row_frame = row_box.text_frame
            row_frame.word_wrap = True

            p = row_frame.paragraphs[0]
            p.text = f"  {supplier['name'][:25]:<25}  {supplier['level']:<10}  {supplier['advantage'][:35]}"
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(51, 51, 51)

            current_top += row_height

        current_top += Inches(0.1)

    return slide

def create_package_solutions_slide(prs, package_images=None):
    """创建方案落地组合包页（第13页）"""
    slide_layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)

    add_title_with_style(slide, "方案落地组合包推荐", Inches(0.4), 28)

    # 三列布局
    col_width = Inches(3)
    col_left = [Inches(0.5), Inches(3.5), Inches(6.5)]
    card_top = Inches(1.3)
    card_height = Inches(5.2)

    colors = [
        RGBColor(102, 153, 204),
        RGBColor(153, 102, 204),
        RGBColor(204, 102, 102)
    ]

    for i, package in enumerate(PACKAGE_SOLUTIONS):
        left = col_left[i]
        current_top = card_top

        # 如果有图片，添加图片
        if package_images and i < len(package_images) and package_images[i] and os.path.exists(package_images[i]):
            try:
                pic_height = Inches(1.5)
                slide.shapes.add_picture(package_images[i], left, current_top, width=col_width, height=pic_height)
                current_top += pic_height + Inches(0.1)
                card_height = Inches(3.6)
            except Exception as e:
                print(f"[WARNING] 添加方案图片失败: {e}")

        # 方案卡片
        card_box = slide.shapes.add_textbox(left, current_top, col_width, card_height)
        card_frame = card_box.text_frame
        card_frame.word_wrap = True

        # 方案名称
        p = card_frame.paragraphs[0]
        p.text = package["name"]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors[i]
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

        # 预算
        p = card_frame.add_paragraph()
        p.text = f"预算：{package['budget']}"
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

        # 描述
        p = card_frame.add_paragraph()
        p.text = package["description"]
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(102, 102, 102)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

        # 核心配置
        p = card_frame.add_paragraph()
        p.text = "核心配置"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(4)

        for item in package["items"]:
            p = card_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(9)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.space_after = Pt(2)

        # 适用对象
        p = card_frame.add_paragraph()
        p.text = f"\n适用：{package['target']}"
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = RGBColor(153, 153, 153)

    return slide

def create_appendix_slide(prs, title, content_items):
    """创建附录页"""
    slide_layout = get_blank_layout(prs)
    slide = prs.slides.add_slide(slide_layout)

    add_title_with_style(slide, title, Inches(0.4), 26)

    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8.8), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    for item in content_items:
        p = content_frame.add_paragraph() if content_frame.paragraphs[0].text else content_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(6)

    return slide

def generate_kitchen_ppt(template_path=None, output_path=None, download_images=True):
    """
    生成现代幸福厨房PPT

    Parameters:
    - template_path: str, PPT模板路径
    - output_path: str, 输出路径
    - download_images: bool, 是否下载图片

    Returns:
    - dict: 包含成功状态和输出文件路径
    """
    try:
        # 使用绝对路径指向项目根目录的Test_result
        script_dir = Path(__file__).parent
        project_root = script_dir.parent.parent.parent  # 回退到项目根目录
        test_result_dir = project_root / "Test_result"

        # 设置默认输出路径
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = str(test_result_dir / f"现代幸福厨房_完整版_{timestamp}.pptx")

        # 确保输出目录存在
        os.makedirs(test_result_dir, exist_ok=True)

        # 图片存储目录 - 使用绝对路径指向项目根目录的Test_result
        image_dir = str(test_result_dir / "images")
        os.makedirs(image_dir, exist_ok=True)
        print(f"[INFO] 输出目录: {test_result_dir}")
        print(f"[INFO] 图片存储目录: {image_dir}")

        # 加载模板
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
            print(f"[INFO] 已加载模板: {template_path}")
            # 清空所有幻灯片
            for _ in range(len(prs.slides)):
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
        else:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            print("[INFO] 创建新PPT（未使用模板）")

        print("[INFO] 开始生成16页PPT...")

        # 准备图片
        section_images = {}
        product_images = {}
        package_images_list = []

        if download_images:
            print("[INFO] 开始下载图片...")
            # 下载分区封面图片
            for section_key in ["A", "B", "C", "D"]:
                paths = prepare_section_images(section_key, SECTIONS_DATA[section_key], image_dir)
                if paths:
                    section_images[section_key] = paths[0]

            # 下载产品图片
            for section_key in ["A", "B", "C", "D"]:
                product_images[section_key] = []
                for product in SECTIONS_DATA[section_key]["products"]:
                    paths = prepare_product_images(product["name"], os.path.join(image_dir, f"products_{section_key}"))
                    product_images[section_key].append(paths)

            # 下载方案组合包图片
            for package in PACKAGE_SOLUTIONS:
                paths = prepare_package_images(package["name"], os.path.join(image_dir, "packages"))
                package_images_list.append(paths[0] if paths else None)

        # 第1页：封面
        print("  [1/16] 封面...")
        cover_image = os.path.join(image_dir, "cover.jpg") if download_images else None
        if download_images:
            try:
                url = "https://picsum.photos/1280/720?random=1"
                response = requests.get(url, timeout=15)
                if response.status_code == 200:
                    with open(cover_image, 'wb') as f:
                        f.write(response.content)
            except:
                cover_image = None
        create_cover_slide(prs, cover_image)

        # 第2页：目录与方法
        print("  [2/16] 目录与方法...")
        create_content_outline_slide(prs)

        # 第3页：五大理念总览
        print("  [3/16] 五大理念总览...")
        create_overview_slide(prs)

        # 第4-5页：A区
        print("  [4-5/16] A区：整体空间与一体化集成...")
        section_a_image = section_images.get("A") if download_images else None
        create_section_cover_slide(prs, "A", SECTIONS_DATA["A"], section_a_image)
        create_products_slide(prs, "A", SECTIONS_DATA["A"]["products"], product_images.get("A"))

        # 第6-7页：B区
        print("  [6-7/16] B区：智能烹饪区...")
        section_b_image = section_images.get("B") if download_images else None
        create_section_cover_slide(prs, "B", SECTIONS_DATA["B"], section_b_image)
        create_products_slide(prs, "B", SECTIONS_DATA["B"]["products"], product_images.get("B"))

        # 第8-9页：C区
        print("  [8-9/16] C区：清洁收纳区...")
        section_c_image = section_images.get("C") if download_images else None
        create_section_cover_slide(prs, "C", SECTIONS_DATA["C"], section_c_image)
        create_products_slide(prs, "C", SECTIONS_DATA["C"]["products"], product_images.get("C"))

        # 第10-11页：D区
        print("  [10-11/16] D区：环境与材料...")
        section_d_image = section_images.get("D") if download_images else None
        create_section_cover_slide(prs, "D", SECTIONS_DATA["D"], section_d_image)
        create_products_slide(prs, "D", SECTIONS_DATA["D"]["products"], product_images.get("D"))

        # 第12页：供应商资源矩阵
        print("  [12/16] 供应商资源矩阵...")
        create_supplier_matrix_slide(prs)

        # 第13页：方案落地组合包
        print("  [13/16] 方案落地组合包...")
        create_package_solutions_slide(prs, package_images_list if download_images else None)

        # 第14页：数据来源
        print("  [14/16] 数据来源...")
        sources_content = [
            "【数据来源】",
            "",
            "企业官方：老板电器、方太、华帝、欧派、索菲亚等官网及新闻稿",
            "行业展会：2024AWE中国家电及消费电子博览会、KBIS2024、IFA2024",
            "权威媒体：新浪地产、腾讯新闻、新华网、36氪等",
            "行业报告：2024中国家电零售与创新白皮书、厨电行业报告",
            "专业机构：中国家用电器协会、iF设计奖",
            "",
            "【重要说明】",
            "• 本PPT所列产品信息均基于2024年公开资料整理",
            "• 具体产品参数、价格、供货周期请以供应商最新报价为准",
            "• 建议在采购前进行供应商实地考察与样品确认",
            "• 图片来源于免费图库，仅供参考，实际产品以供应商提供为准"
        ]
        create_appendix_slide(prs, "数据来源与说明", sources_content)

        # 第15页：免责声明
        print("  [15/16] 免责声明...")
        disclaimer_content = [
            "【免责声明】",
            "",
            "1. 本报告仅供参考，不构成任何采购承诺或法律约束。",
            "",
            "2. 实际采购决策请结合项目具体情况进行综合评估，包括但不限于：",
            "   • 预算约束与成本效益分析",
            "   • 供应商交付能力与售后服务",
            "   • 产品适配性与工程条件",
            "   • 合同条款与风险管控",
            "",
            "3. 所有产品信息、技术参数、价格等均以供应商官方最新发布为准。",
            "",
            "4. 建议在采购前进行多方案比选和样品确认。",
            "",
            f"编制单位：上海金地采购部",
            f"编制日期：{datetime.now().strftime('%Y年%m月%d日')}"
        ]
        create_appendix_slide(prs, "免责声明", disclaimer_content)

        # 第16页：完整来源列表
        print("  [16/16] 完整来源列表...")
        refs_content = [
            "【完整来源列表】",
            "",
            "【A区 - 整体空间与一体化集成】",
            "• 欧派家居：www.oppein.com | 2024整家定制2.0 | 欧派2023可持续发展报告",
            "• 智小金MetaBox：2022年发布 | 装配式建筑技术规程",
            "• 可变隔断系统：2024大家居材艺趋势白皮书 | 2025年温岭市幼儿园应用案例",
            "",
            "【B区 - 智能烹饪区】",
            "• 老板电器食神AI：2024AWE | 2024年营收公告 | 2025年食神大模型发布",
            "• 方太集成烹饪中心：2019年首创 | 2025厨电三强报告 | AWE2024",
            "• 华帝AI智慧烹饪中心：2024产品发布会 | AWE2025报道",
            "• COLMO TURING 2.0：IFA2024 | 美的集团2024半年报",
            "",
            "【C区 - 清洁收纳区】",
            "• 方太水槽洗碗机：KBIS2024 | 2024年度质量诚信报告",
            "• 海尔双面洗：AWE2024 | 2024洗碗机选购指南 | 2024洗碗机行业报告",
            "• 卡萨帝抽屉式洗碗机：iF设计奖 | IFA2024 | 2024家电零售白皮书",
            "",
            "【D区 - 环境与材料】",
            "• 磐珉零硅石英石：2024石英石产品目录 | 可持续发展报告",
            "• 威洋无机水磨石：2025新产品发布 | 威洋高性能无机水磨石介绍",
            "• VIATERA石英石：LX Hausys官网产品页面",
            "• 欧派环保板材：2023可持续发展报告 | 绿色家居产业链"
        ]
        create_appendix_slide(prs, "完整来源列表", refs_content)

        # 保存PPT
        print(f"[INFO] 正在保存PPT到: {output_path}")
        prs.save(output_path)

        print(f"[SUCCESS] PPT生成成功！")
        print(f"[INFO] 输出文件: {output_path}")
        print(f"[INFO] 总页数: {len(prs.slides)}")

        return {
            "success": True,
            "output_file": output_path,
            "slide_count": len(prs.slides)
        }

    except Exception as e:
        print(f"[ERROR] PPT生成失败: {str(e)}")
        import traceback
        traceback.print_exc()
        return {
            "success": False,
            "error": str(e)
        }

if __name__ == "__main__":
    # 指定模板路径和输出路径
    script_dir = Path(__file__).parent
    template_path = str(script_dir / ".." / ".." / ".." / "Test_result" / "PPT主题模板.pptx")
    output_path = str(script_dir / ".." / ".." / ".." / "Test_result" / "现代幸福厨房_最终版.pptx")

    # 生成PPT
    result = generate_kitchen_ppt(
        template_path=template_path,
        output_path=output_path,
        download_images=True
    )

    if result["success"]:
        print("\n" + "="*60)
        print("现代幸福厨房PPT生成完成！")
        print("="*60)
        print(f"输出文件: {result['output_file']}")
        print(f"总页数: {result['slide_count']}")
        print("\n页面结构：")
        print("  第1页：封面")
        print("  第2页：目录与方法")
        print("  第3页：五大理念总览")
        print("  第4-5页：A区（整体空间与一体化集成）")
        print("  第6-7页：B区（智能烹饪区）")
        print("  第8-9页：C区（清洁收纳区）")
        print("  第10-11页：D区（环境与材料）")
        print("  第12页：供应商资源矩阵")
        print("  第13页：方案落地组合包")
        print("  第14-16页：附录（数据来源、免责声明、完整来源列表）")
        print("="*60)
    else:
        print("\nPPT生成失败，请检查错误信息。")
