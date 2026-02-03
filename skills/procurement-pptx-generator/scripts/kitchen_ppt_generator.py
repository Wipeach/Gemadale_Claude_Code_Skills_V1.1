#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
现代幸福厨房 PPT 生成器
根据提示词生成包含2024+创新产品与供应链推荐的PPT
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from datetime import datetime
from pathlib import Path
from typing import Dict, List, Any

# ==================== 数据定义 ====================

# 幸福厨房五大理念标签
HAPPY_KITCHEN_TAGS = [
    "①一体化集成",
    "②智能场景",
    "③健康安全",
    "④灵活可变",
    "⑤美学永续"
]

# 四大分区数据
SECTIONS_DATA = {
    "A": {
        "title": "A. 整体空间与一体化集成",
        "subtitle": "天地门墙柜燃电一体化、模块化、装配式、系统集成",
        "trends": [
            "门墙柜一体化整装定制成为主流",
            "装配式模块化厨房实现快速搭建",
            "可变空间设计实现开放/封闭切换"
        ],
        "products": [
            {
                "name": "欧派门墙柜一体化整家定制",
                "position": "全屋一体化空间解决方案",
                "tags": ["①一体化集成", "⑤美学永续"],
                "innovation": "通过门墙柜整体配设计，实现厨房、衣柜、浴室、门窗等家具的统一协调；门墙柜自制交付优势",
                "stage": "已上市",
                "params": "覆盖全屋8大空间，客单值超5万",
                "image": "欧派门墙柜一体化",
                "reason": "一体化设计确保风格统一，降低沟通成本，提升交付效率"
            },
            {
                "name": "智小金 MetaBox 模块化厨房",
                "position": "集模块化、装配式为一体的智能收纳系统",
                "tags": ["①一体化集成", "④灵活可变"],
                "innovation": "能与厨柜/衣柜/墙板完美嵌合，实现科技融于家居；装配式设计支持快速安装与拆卸",
                "stage": "已上市",
                "params": "模块化设计，支持多种组合方式",
                "image": "智小金MetaBox",
                "reason": "模块化设计让厨房像乐高一样灵活组合，满足不同户型需求"
            },
            {
                "name": "可变隔断系统（折叠门/滑移岛台）",
                "position": "开放与封闭模式自由切换的空间方案",
                "tags": ["④灵活可变", "①一体化集成"],
                "innovation": "折叠门、滑移岛台使厨房可在开放共享与封闭专注两种模式间切换；无下轨设计便于清洁",
                "stage": "已量产",
                "params": "四联动互推拉、无下轨设计",
                "image": "可变隔断系统",
                "reason": "解决中式烹饪油烟问题，同时保持社交属性"
            }
        ],
        "suppliers": [
            {"name": "欧派家居", "advantage": "上市公司，全屋定制龙头，门墙柜一体化交付能力强", "level": "国内高端"},
            {"name": "索菲亚", "advantage": "定制家居领军品牌，整家定制解决方案完善", "level": "国内一线"},
            {"name": "志邦家居", "advantage": "橱衣双品牌统一，全屋定制经验丰富", "level": "国内一线"},
            {"name": "金牌厨柜", "advantage": "专注于厨柜领域，工程案例丰富", "level": "国内一线"},
            {"name": "海尔全屋家居", "advantage": "家电与家居深度融合优势", "level": "国内知名"}
        ]
    },
    "B": {
        "title": "B. 智能烹饪区",
        "subtitle": "AI/IoT联动：灶烟蒸烤冰箱、AI菜谱与流程自动化、视觉识别",
        "trends": [
            "AI烹饪大模型实现个性化菜谱推荐",
            "烟灶蒸烤一体化联动控制",
            "视觉识别技术实现食材智能检测"
        ],
        "products": [
            {
                "name": "老板电器 AI烹饪大模型「食神」",
                "position": "全球首个烹饪垂直领域AI大模型",
                "tags": ["②智能场景", "①一体化集成"],
                "innovation": "AI烟灶联动系统，自动火候控制；支持200+智能菜谱；全链路烹饪解决方案",
                "stage": "已上市",
                "params": "双平台（小程序/APP），2024年已商业闭环",
                "image": "老板食神AI",
                "reason": "AI让烹饪更简单，新手也能做出大师级菜品"
            },
            {
                "name": "方太集成烹饪中心",
                "position": "烟灶蒸烤一体化的高端烹饪解决方案",
                "tags": ["①一体化集成", "②智能场景"],
                "innovation": "一体集成油烟机、燃气灶、蒸烤箱等多种功能；上排集成设计节省空间；AI智能菜谱",
                "stage": "已上市",
                "params": "10㎡以上厨房场景，套系销售占比63%",
                "image": "方太集成烹饪中心",
                "reason": "一机多用，完美解决小户型厨房空间不足问题"
            },
            {
                "name": "华帝 AI智慧集成烹饪中心",
                "position": "集成烟灶区域与蒸烤区域的智能烹饪系统",
                "tags": ["②智能场景", "①一体化集成"],
                "innovation": "双腔蒸烤一体机，可同时进行蒸和烤；AI联动控制；跨设备智能协同",
                "stage": "已上市",
                "params": "双腔设计，支持多种烹饪模式组合",
                "image": "华帝AI智慧烹饪中心",
                "reason": "双腔设计大幅提升烹饪效率，蒸烤同时进行"
            }
        ],
        "suppliers": [
            {"name": "老板电器", "advantage": "厨电行业龙头，AI烹饪技术领先，工程案例丰富", "level": "国内高端"},
            {"name": "方太", "advantage": "高端厨电领导品牌，集成烹饪中心开创者", "level": "国内高端"},
            {"name": "华帝", "advantage": "厨电三强之一，智能化转型成果显著", "level": "国内一线"},
            {"name": "COLMO", "advantage": "美的集团高端品牌，AI科技家电标杆", "level": "国际高端"},
            {"name": "博世家电", "advantage": "德系精工，嵌入式厨电技术领先", "level": "国际高端"}
        ]
    },
    "C": {
        "title": "C. 清洁收纳区",
        "subtitle": "洗碗/净水/水槽/垃圾处理/收纳五金与智能整理",
        "trends": [
            "水槽式洗碗机成为中式厨房新选择",
            "AI双面洗技术解决中餐重油污问题",
            "抗菌水槽与智能净水一体化"
        ],
        "products": [
            {
                "name": "方太水槽洗碗机（高能气泡洗）",
                "position": "专为中餐设计的三合一水槽洗碗机",
                "tags": ["③健康安全", "①一体化集成"],
                "innovation": "高能气泡洗技术，更适合中国家庭；水槽与洗碗机融合，节省空间；一机三用（水槽+洗碗+净洗）",
                "stage": "已上市",
                "params": "净水流量2.1L/min，额定功率98W",
                "image": "方太水槽洗碗机",
                "reason": "完美适配中式厨房，解决重油污清洗难题"
            },
            {
                "name": "海尔AI双面洗洗碗机",
                "position": "专利H型中喷淋臂双面洁净技术",
                "tags": ["③健康安全", "②智能场景"],
                "innovation": "H型中喷淋臂实现上下双面包裹式冲刷；56000Pa水压+80度蒸汽变温技术",
                "stage": "已上市",
                "params": "双面洗技术，洗净率提升30%",
                "image": "海尔双面洗洗碗机",
                "reason": "双面洗技术彻底解决中餐重油污清洗痛点"
            },
            {
                "name": "卡萨帝抽屉式洗碗机",
                "position": "液力悬浮喷淋臂直连，无水压损耗",
                "tags": ["③健康安全", "⑤美学永续"],
                "innovation": "抽屉式设计符合人体工学；液力悬浮喷淋臂直连电机，水流强度更有保证；节能与清洗效率兼顾",
                "stage": "已上市",
                "params": "大16套碗盘同时洗，变频电机节能高效",
                "image": "卡萨帝抽屉式洗碗机",
                "reason": "抽屉式设计优雅便捷，用户体验极佳"
            }
        ],
        "suppliers": [
            {"name": "方太", "advantage": "水槽洗碗机品类开创者，技术专利丰富", "level": "国内高端"},
            {"name": "海尔智家", "advantage": "双面洗技术专利，洗碗机市场份额领先", "level": "国内一线"},
            {"name": "西门子家电", "advantage": "洗碗机技术领先，品质可靠", "level": "国际高端"},
            {"name": "卡萨帝", "advantage": "海尔高端品牌，抽屉式洗碗机创新者", "level": "国内高端"},
            {"name": "美的", "advantage": "性价比高，产品线丰富", "level": "国内一线"}
        ]
    },
    "D": {
        "title": "D. 环境与材料",
        "subtitle": "台面、柜体板材、墙地面、防滑抗菌易清洁、空气治理与低VOC材料",
        "trends": [
            "零硅板材等健康环保材料成为趋势",
            "石英石/岩板台面耐污抗指纹",
            "抗菌材料在厨房空间广泛应用"
        ],
        "products": [
            {
                "name": "磐珉零硅石英石台面",
                "position": "环保健康的新型石英石材料",
                "tags": ["③健康安全", "⑤美学永续"],
                "innovation": "零硅板材助于营造更健康的工作环境；深度融合循环材料和矿物；符合全球可持续发展标准",
                "stage": "已上市",
                "params": "低VOC、环保等级E0级以上",
                "image": "磐珉零硅石英石",
                "reason": "健康环保，符合绿色建筑与可持续发展要求"
            },
            {
                "name": "威洋高性能无机水磨石/石英石",
                "position": "纯无机基因的高端装饰材料",
                "tags": ["③健康安全", "⑤美学永续"],
                "innovation": "含有天然石英质类材料；常用于厨房卫生间台面板、墙面地面铺贴；耐污抗指纹",
                "stage": "已上市",
                "params": "耐污、耐高温、易清洁",
                "image": "威洋无机水磨石",
                "reason": "耐久性强，维护成本低，适合高频使用场景"
            },
            {
                "name": "欧派低VOC环保柜体板材",
                "position": "绿色环保的全屋定制板材",
                "tags": ["③健康安全", "⑤美学永续"],
                "innovation": "系统构建绿色家居产业链；产品绿色全生命周期管理；废气、噪声排放全部达标",
                "stage": "已上市",
                "params": "环保等级E0级，符合国家室内装饰装修材料标准",
                "image": "欧派环保板材",
                "reason": "从源头控制室内空气污染，守护家人健康"
            }
        ],
        "suppliers": [
            {"name": "磐珉新材", "advantage": "石英石专业供应商，零硅板材技术领先", "level": "国内知名"},
            {"name": "威洋石材", "advantage": "石英石和人造石专业制造商，工程经验丰富", "level": "国内知名"},
            {"name": "LX Hausys（VIATERA）", "advantage": "韩系品牌，石英石台面质量优异", "level": "国际高端"},
            {"name": "欧派家居", "advantage": "环保板材供应链完善，全生命周期管理", "level": "国内高端"},
            {"name": "兔宝宝", "advantage": "环保板材龙头，E0/ENF级标准引领者", "level": "国内知名"}
        ]
    }
}

# 方案落地组合包
PACKAGE_SOLUTIONS = [
    {
        "name": "轻量升级包",
        "description": "基础智能化改造，适合存量房升级",
        "items": [
            "基础智能烟灶联动套装",
            "水槽式洗碗机（方太/海尔）",
            "普通石英石台面",
            "基础收纳五金升级"
        ],
        "budget": "3-5万",
        "target": "存量房改造、预算有限项目"
    },
    {
        "name": "中配智能包",
        "description": "智能化与一体化并重，适合新建项目",
        "items": [
            "集成烹饪中心（方太/老板）",
            "嵌入式洗碗机（海尔/西门子）",
            "标准石英石台面+环保板材",
            "可移门隔断系统",
            "净水设备（科勒/方太）"
        ],
        "budget": "8-12万",
        "target": "新建精装房、改善型项目"
    },
    {
        "name": "旗舰一体化包",
        "description": "全屋一体化+AI智能，高端定位",
        "items": [
            "门墙柜一体化定制（欧派/索菲亚）",
            "AI烹饪大系统（老板食神+集成烹饪中心）",
            "高端洗碗机（卡萨帝/西门子）",
            "零硅石英石/岩板台面",
            "全屋智能净水系统",
            "可变空间系统（折叠门+滑移岛台）"
        ],
        "budget": "15-25万",
        "target": "高端豪宅、旗舰项目"
    }
]

# ==================== PPT生成核心函数 ====================

def create_title_slide(prs, template_slide_layout=None):
    """创建封面页"""
    if template_slide_layout:
        slide = prs.slides.add_slide(template_slide_layout)
    else:
        blank_layout = prs.slide_layouts[6]  # 空白布局
        slide = prs.slides.add_slide(blank_layout)

    # 添加主标题
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)

    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    p = title_frame.paragraphs[0]
    p.text = "现代幸福厨房"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.alignment = PP_ALIGN.CENTER

    # 添加副标题
    subtitle_top = Inches(3.8)
    subtitle_box = slide.shapes.add_textbox(left, subtitle_top, width, Inches(0.8))
    subtitle_frame = subtitle_box.text_frame

    p = subtitle_frame.paragraphs[0]
    p.text = "2024+ 创新产品与供应链推荐"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(102, 102, 102)
    p.alignment = PP_ALIGN.CENTER

    # 添加部门信息
    dept_top = Inches(5.2)
    dept_box = slide.shapes.add_textbox(left, dept_top, width, Inches(0.6))
    dept_frame = dept_box.text_frame

    p = dept_frame.paragraphs[0]
    p.text = "上海金地采购部"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(153, 153, 153)
    p.alignment = PP_ALIGN.CENTER

    # 添加日期
    date_top = Inches(5.8)
    date_box = slide.shapes.add_textbox(left, date_top, width, Inches(0.5))
    date_frame = date_box.text_frame

    p = date_frame.paragraphs[0]
    p.text = datetime.now().strftime("%Y年%m月")
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(153, 153, 153)
    p.alignment = PP_ALIGN.CENTER

    return slide


def create_content_slide(prs, template_slide_layout=None):
    """创建目录页"""
    if template_slide_layout:
        slide = prs.slides.add_slide(template_slide_layout)
    else:
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "目录"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)

    # 目录内容
    content_left = Inches(1)
    content_top = Inches(1.8)
    content_width = Inches(8)

    content_box = slide.shapes.add_textbox(content_left, content_top, content_width, Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    contents = [
        "01 筛选口径与方法论",
        "02 幸福厨房五大理念与一体化架构",
        "03 A区：整体空间与一体化集成",
        "04 B区：智能烹饪区",
        "05 C区：清洁收纳区",
        "06 D区：环境与材料",
        "07 供应商资源矩阵",
        "08 方案落地组合包推荐",
        "09 数据来源与附录"
    ]

    for i, item in enumerate(contents):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()

        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(12)

    return slide


def create_methodology_slide(prs, template_slide_layout=None):
    """创建筛选口径与方法论页"""
    if template_slide_layout:
        slide = prs.slides.add_slide(template_slide_layout)
    else:
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "筛选口径与方法论"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)

    # 筛选标准
    criteria_left = Inches(0.8)
    criteria_top = Inches(1.6)
    criteria_box = slide.shapes.add_textbox(criteria_left, criteria_top, Inches(4), Inches(5))
    criteria_frame = criteria_box.text_frame
    criteria_frame.word_wrap = True

    p = criteria_frame.paragraphs[0]
    p.text = "筛选标准"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(12)

    criteria_list = [
        "时间范围：2024年至今发布/量产产品",
        "产品阶段：已上市 / 刚量产 / 概念阶段",
        "创新性：技术突破/模式创新/材料革新",
        "可落地性：供应链完善、工程案例丰富",
        "数据来源：多源交叉验证（官网/展会/权威媒体）"
    ]

    for item in criteria_list:
        p = criteria_frame.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(8)
        p.level = 0

    # 评价维度
    eval_left = Inches(5.5)
    eval_box = slide.shapes.add_textbox(eval_left, criteria_top, Inches(4), Inches(5))
    eval_frame = eval_box.text_frame
    eval_frame.word_wrap = True

    p = eval_frame.paragraphs[0]
    p.text = "评价维度"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(12)

    eval_list = [
        "技术创新性（传感器/算法/材料）",
        "集成能力（系统联动/空间优化）",
        "健康安全（抗菌/环保/低VOC）",
        "用户体验（易用性/智能化/美观）",
        "供应链稳定性（品牌/产能/交付）"
    ]

    for item in eval_list:
        p = eval_frame.add_paragraph()
        p.text = f"● {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(8)
        p.level = 0

    return slide


def create_overview_slide(prs, template_slide_layout=None):
    """创建幸福厨房五大理念总览页"""
    if template_slide_layout:
        slide = prs.slides.add_slide(template_slide_layout)
    else:
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "幸福厨房五大理念 + 一体化架构"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)

    # 五大理念
    tags_left = Inches(0.8)
    tags_top = Inches(1.4)
    tags_width = Inches(8.5)

    tags_box = slide.shapes.add_textbox(tags_left, tags_top, tags_width, Inches(2.5))
    tags_frame = tags_box.text_frame
    tags_frame.word_wrap = True

    p = tags_frame.paragraphs[0]
    p.text = "幸福厨房五大理念"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(10)

    for tag in HAPPY_KITCHEN_TAGS:
        p = tags_frame.add_paragraph()
        p.text = tag
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(6)

    # 一体化架构说明
    arch_left = Inches(0.8)
    arch_top = Inches(4.2)
    arch_box = slide.shapes.add_textbox(arch_left, arch_top, tags_width, Inches(2.5))
    arch_frame = arch_box.text_frame
    arch_frame.word_wrap = True

    p = arch_frame.paragraphs[0]
    p.text = "天地门墙柜燃电一体化架构"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(10)

    arch_desc = [
        "天：吊顶系统、烟道排风、照明集成",
        "地：地柜、地台、水电气暗埋集成",
        "门：智能移门、折叠门、可变隔断",
        "墙：墙板、背景墙、储物墙一体化",
        "柜：地柜、吊柜、高柜、岛台集成",
        "燃：燃气灶、集成灶、燃气管道安全",
        "电：智能电器、电力布局、智能控制"
    ]

    for item in arch_desc:
        p = arch_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(5)

    return slide


def create_section_slide(prs, section_key, section_data, template_slide_layout=None):
    """创建分区封面页"""
    if template_slide_layout:
        slide = prs.slides.add_slide(template_slide_layout)
    else:
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

    # 分区标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = section_data["title"]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)

    # 分区副标题
    subtitle_top = Inches(1.5)
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), subtitle_top, Inches(9), Inches(0.6))
    subtitle_frame = subtitle_box.text_frame
    p = subtitle_frame.paragraphs[0]
    p.text = section_data["subtitle"]
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(102, 102, 102)

    # 关键趋势
    trends_top = Inches(2.5)
    trends_box = slide.shapes.add_textbox(Inches(0.8), trends_top, Inches(8.5), Inches(4))
    trends_frame = trends_box.text_frame
    trends_frame.word_wrap = True

    p = trends_frame.paragraphs[0]
    p.text = "关键趋势"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(15)

    for i, trend in enumerate(section_data["trends"], 1):
        p = trends_frame.add_paragraph()
        p.text = f"{i}. {trend}"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(12)

    return slide


def create_product_cards_slide(prs, section_key, products, slide_index, template_slide_layout=None):
    """创建产品卡片页（每页2-3个产品）"""
    if template_slide_layout:
        slide = prs.slides.add_slide(template_slide_layout)
    else:
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

    # 页面标题
    title_text = f"{SECTIONS_DATA[section_key]['title']} - 产品推荐"
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)

    # 创建产品卡片
    card_height = 2.8
    card_spacing = 0.2
    start_top = 1.2

    for i, product in enumerate(products):
        card_top = start_top + i * (card_height + card_spacing)

        # 卡片背景框（用文本框模拟）
        card_box = slide.shapes.add_textbox(
            Inches(0.6),
            Inches(card_top),
            Inches(8.8),
            Inches(card_height)
        )
        card_frame = card_box.text_frame
        card_frame.word_wrap = True

        # 产品名称和定位
        p = card_frame.paragraphs[0]
        p.text = f"{product['name']}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(4)

        p = card_frame.add_paragraph()
        p.text = product['position']
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = RGBColor(102, 102, 102)
        p.space_after = Pt(6)

        # 标签
        tags_str = " | ".join(product['tags'])
        p = card_frame.add_paragraph()
        p.text = f"标签：{tags_str}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0, 102, 204)
        p.space_after = Pt(6)

        # 创新点
        p = card_frame.add_paragraph()
        innovation_short = product['innovation'][:80] + "..." if len(product['innovation']) > 80 else product['innovation']
        p.text = f"创新：{innovation_short}"
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(4)

        # 阶段和推荐理由
        p = card_frame.add_paragraph()
        p.text = f"阶段：{product['stage']}  |  理由：{product['reason'][:60]}..."
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(102, 102, 102)

    return slide


def create_supplier_matrix_slide(prs, template_slide_layout=None):
    """创建供应商资源矩阵页"""
    if template_slide_layout:
        slide = prs.slides.add_slide(template_slide_layout)
    else:
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "供应商资源矩阵"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)

    # 创建供应商表格（用文本框模拟）
    table_left = Inches(0.5)
    table_top = Inches(1.2)
    row_height = 0.5

    # 表头
    header_box = slide.shapes.add_textbox(table_left, table_top, Inches(9), Inches(row_height))
    header_frame = header_box.text_frame
    header_frame.word_wrap = True

    p = header_frame.paragraphs[0]
    p.text = "品类          供应商名称                                      级别        核心优势"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # 供应商数据（整合四个分区）
    all_suppliers = []
    for section_key, section_data in SECTIONS_DATA.items():
        for supplier in section_data["suppliers"]:
            all_suppliers.append({
                "category": section_key,
                "name": supplier["name"],
                "level": supplier["level"],
                "advantage": supplier["advantage"]
            })

    # 分组展示（按A/B/C/D分区）
    current_section = None
    current_top = table_top + Inches(row_height + 0.1)

    for supplier in all_suppliers:
        if supplier["category"] != current_section:
            # 新分区标题
            current_section = supplier["category"]
            section_box = slide.shapes.add_textbox(table_left, current_top, Inches(9), Inches(0.4))
            section_frame = section_box.text_frame
            p = section_frame.paragraphs[0]
            p.text = f"{SECTIONS_DATA[current_section]['title']}"
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = RGBColor(51, 51, 51)
            current_top += Inches(0.5)

        # 供应商行
        row_box = slide.shapes.add_textbox(table_left, current_top, Inches(9), Inches(row_height))
        row_frame = row_box.text_frame
        row_frame.word_wrap = True

        p = row_frame.paragraphs[0]
        p.text = f"          {supplier['name'][:30]:<30}  {supplier['level']:<12}  {supplier['advantage'][:40]}"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(51, 51, 51)

        current_top += Inches(row_height + 0.05)

    return slide


def create_package_solutions_slide(prs, template_slide_layout=None):
    """创建方案落地组合包页"""
    if template_slide_layout:
        slide = prs.slides.add_slide(template_slide_layout)
    else:
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "方案落地组合包推荐"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)

    # 三列布局展示三套方案
    col_width = 3
    col_left = [Inches(0.5), Inches(3.5), Inches(6.5)]
    card_top = Inches(1.2)
    card_height = Inches(4.5)

    colors = [
        RGBColor(102, 153, 204),   # 蓝色
        RGBColor(153, 102, 204),   # 紫色
        RGBColor(204, 102, 102)    # 红色
    ]

    for i, package in enumerate(PACKAGE_SOLUTIONS):
        left = col_left[i]

        # 方案卡片
        card_box = slide.shapes.add_textbox(left, card_top, Inches(col_width), card_height)
        card_frame = card_box.text_frame
        card_frame.word_wrap = True

        # 方案名称
        p = card_frame.paragraphs[0]
        p.text = package["name"]
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = colors[i]
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(8)

        # 预算范围
        p = card_frame.add_paragraph()
        p.text = f"预算：{package['budget']}"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(10)

        # 描述
        p = card_frame.add_paragraph()
        p.text = package["description"]
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(102, 102, 102)
        p.space_after = Pt(10)

        # 核心配置
        p = card_frame.add_paragraph()
        p.text = "核心配置："
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(5)

        for item in package["items"]:
            p = card_frame.add_paragraph()
            p.text = f"• {item}"
            p.font.size = Pt(10)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.space_after = Pt(3)
            p.level = 0

        # 适用对象
        p = card_frame.add_paragraph()
        p.text = f"\n适用：{package['target']}"
        p.font.size = Pt(10)
        p.font.italic = True
        p.font.color.rgb = RGBColor(153, 153, 153)

    return slide


def create_appendix_slide(prs, template_slide_layout=None):
    """创建数据来源与附录页"""
    if template_slide_layout:
        slide = prs.slides.add_slide(template_slide_layout)
    else:
        blank_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_layout)

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    title_frame = title_box.text_frame
    p = title_frame.paragraphs[0]
    p.text = "数据来源与附录"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)

    # 数据来源列表
    sources_left = Inches(0.6)
    sources_top = Inches(1.2)
    sources_box = slide.shapes.add_textbox(sources_left, sources_top, Inches(8.8), Inches(5.5))
    sources_frame = sources_box.text_frame
    sources_frame.word_wrap = True

    p = sources_frame.paragraphs[0]
    p.text = "数据来源"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = RGBColor(51, 51, 51)
    p.space_after = Pt(10)

    sources = [
        "【企业官方】老板电器、方太、华帝、欧派、索菲亚等官网及新闻稿",
        "【行业展会】2024中国家电及消费电子博览会（AWE）、中国厨卫展",
        "【权威媒体】新浪地产、腾讯新闻、新华网家电频道等",
        "【行业报告】2024中国家电零售与创新白皮书、厨电行业报告",
        "【行业协会】中国家用电器协会、中国建筑卫生陶瓷协会",
        "",
        "重要说明：",
        "• 本PPT所列产品信息均基于2024年公开资料整理",
        "• 具体产品参数、价格、供货周期请以供应商最新报价为准",
        "• 建议在采购前进行供应商实地考察与样品确认"
    ]

    for source in sources:
        p = sources_frame.add_paragraph()
        p.text = source
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(51, 51, 51)
        p.space_after = Pt(6)

    # 添加免责声明
    disclaimer_top = Inches(5.5)
    disclaimer_box = slide.shapes.add_textbox(Inches(0.6), disclaimer_top, Inches(8.8), Inches(0.8))
    disclaimer_frame = disclaimer_box.text_frame

    p = disclaimer_frame.paragraphs[0]
    p.text = "免责声明：本报告仅供参考，不构成任何采购承诺或法律约束。实际采购决策请结合项目具体情况进行综合评估。"
    p.font.size = Pt(9)
    p.font.italic = True
    p.font.color.rgb = RGBColor(153, 153, 153)

    return slide


def generate_kitchen_ppt(template_path=None, output_path=None):
    """
    生成现代幸福厨房PPT

    Parameters:
    - template_path: str, PPT模板路径（可选）
    - output_path: str, 输出路径（可选，默认为Test_result目录）

    Returns:
    - dict: 包含成功状态和输出文件路径
    """
    try:
        # 设置默认输出路径
        if output_path is None:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            output_path = str(Path("Test_result") / f"现代幸福厨房_{timestamp}.pptx")

        # 确保输出目录存在
        os.makedirs(os.path.dirname(output_path) if os.path.dirname(output_path) else ".", exist_ok=True)

        # 加载模板或创建新PPT
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
            print(f"[INFO] 已加载模板: {template_path}")
        else:
            prs = Presentation()
            prs.slide_width = Inches(10)
            prs.slide_height = Inches(7.5)
            print("[INFO] 使用空白模板创建新PPT")

        # 清空模板中的所有幻灯片（如果是从模板加载）
        if template_path and os.path.exists(template_path):
            # 删除所有现有幻灯片
            for _ in range(len(prs.slides)):
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]

        # 获取空白布局
        blank_layout = prs.slide_layouts[6]

        # ==================== 生成各页 ====================
        print("[INFO] 开始生成PPT页面...")

        # 第1页：封面
        print("  [1/16] 生成封面...")
        create_title_slide(prs, blank_layout)

        # 第2页：目录
        print("  [2/16] 生成目录...")
        create_content_slide(prs, blank_layout)

        # 第3页：筛选口径与方法论
        print("  [3/16] 生成筛选口径与方法论...")
        create_methodology_slide(prs, blank_layout)

        # 第4页：五大理念总览
        print("  [4/16] 生成五大理念总览...")
        create_overview_slide(prs, blank_layout)

        # 第5-6页：A区（整体空间与一体化集成）
        print("  [5-6/16] 生成A区内容...")
        create_section_slide(prs, "A", SECTIONS_DATA["A"], blank_layout)
        create_product_cards_slide(prs, "A", SECTIONS_DATA["A"]["products"], 1, blank_layout)

        # 第7-8页：B区（智能烹饪区）
        print("  [7-8/16] 生成B区内容...")
        create_section_slide(prs, "B", SECTIONS_DATA["B"], blank_layout)
        create_product_cards_slide(prs, "B", SECTIONS_DATA["B"]["products"], 1, blank_layout)

        # 第9-10页：C区（清洁收纳区）
        print("  [9-10/16] 生成C区内容...")
        create_section_slide(prs, "C", SECTIONS_DATA["C"], blank_layout)
        create_product_cards_slide(prs, "C", SECTIONS_DATA["C"]["products"], 1, blank_layout)

        # 第11-12页：D区（环境与材料）
        print("  [11-12/16] 生成D区内容...")
        create_section_slide(prs, "D", SECTIONS_DATA["D"], blank_layout)
        create_product_cards_slide(prs, "D", SECTIONS_DATA["D"]["products"], 1, blank_layout)

        # 第13页：供应商资源矩阵
        print("  [13/16] 生成供应商资源矩阵...")
        create_supplier_matrix_slide(prs, blank_layout)

        # 第14-15页：方案落地组合包
        print("  [14/16] 生成方案落地组合包...")
        create_package_solutions_slide(prs, blank_layout)

        # 第16页：数据来源与附录
        print("  [16/16] 生成数据来源与附录...")
        create_appendix_slide(prs, blank_layout)

        # 保存PPT
        print(f"[INFO] 正在保存PPT到: {output_path}")
        prs.save(output_path)

        print(f"[SUCCESS] PPT生成成功！")
        print(f"[INFO] 输出文件: {output_path}")
        print(f"[INFO] 总页数: {len(prs.slides)}")

        return {
            "success": True,
            "output_file": output_path,
            "slide_count": len(prs.slides)
        }

    except Exception as e:
        print(f"[ERROR] PPT生成失败: {str(e)}")
        import traceback
        traceback.print_exc()
        return {
            "success": False,
            "error": str(e)
        }


if __name__ == "__main__":
    # 运行PPT生成
    result = generate_kitchen_ppt()

    if result["success"]:
        print("\n" + "="*50)
        print("现代幸福厨房PPT生成完成！")
        print("="*50)
        print(f"输出文件: {result['output_file']}")
        print(f"总页数: {result['slide_count']}")
        print("="*50)
    else:
        print("\nPPT生成失败，请检查错误信息。")
