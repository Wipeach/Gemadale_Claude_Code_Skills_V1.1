#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Insert text from a text file into a PPT slide (page 4) as a table.
Parses renovation details and adds a table to the bottom-right quarter of the slide.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from datetime import datetime

def run(project_name):
    """
    从文本文件中提取装修详情，插入到 PowerPoint 第四页右下角，占约1/4内容。
    
    Args:
        project_name (str): 项目名称 (例如, '华发四季半岛')。
    
    Returns:
        dict: 包含操作结果的字典。
    """
    timestamp = datetime.now().strftime("%Y%m%d")
    input_file = f"resources/working_data/{project_name}_{timestamp}/{project_name}_基本信息.txt"
    pptx_path = f"resources/working_data/{project_name}_{timestamp}/processed_data/{project_name}_gemdale_housing_project_template.pptx"
    
    renovation_data = []
    current_section = None
    
    try:
        with open(input_file, 'r', encoding='utf-8') as file:
            lines = file.readlines()
            for line in lines:
                line = line.strip()
                if line == "装修情况:":
                    current_section = "renovation"
                    continue
                if current_section == "renovation":
                    if line == "预证信息:":  # 遇到下一节停止
                        break
                    renovation_data.append(line)
    except FileNotFoundError:
        print(f"错误：文件 {input_file} 未找到。")
        return {"status": "error", "message": f"文件 {input_file} 未找到"}
    except Exception as e:
        print(f"读取文件时出错：{e}")
        return {"status": "error", "message": f"读取文件时出错：{e}"}

    print("原始装修数据：")
    for line in renovation_data:
        print(f"  {line}")

    table_data = []
    headers = ["房间类型", "地面", "墙面", "配置"]
    current_room = None
    room_info = {}
    pending_key = None

    for line in renovation_data:
        line = line.strip()
        if not line:
            continue
        if line in ["室内装修", "全装修", "卧室", "厨房", "客厅", "卫生间"]:
            if current_room and room_info:
                table_data.append([
                    current_room,
                    room_info.get("地面", "-"),
                    room_info.get("墙面", "-"),
                    room_info.get("配置", "-")
                ])
                room_info = {}
            current_room = line
            pending_key = None
        elif line in ["装修价格", "外立面风格", "园林风格"]:
            pending_key = line
        elif line.endswith(":"):
            pending_key = line[:-1]
        elif pending_key:
            if pending_key in ["地面", "墙面", "配置"] and current_room:
                room_info[pending_key] = line
            elif pending_key in ["装修价格", "外立面风格", "园林风格"]:
                table_data.append([pending_key, line, "-", "-"])
            pending_key = None

    if current_room and room_info:
        table_data.append([
            current_room,
            room_info.get("地面", "-"),
            room_info.get("墙面", "-"),
            room_info.get("配置", "-")
        ])

    print("解析后的表格数据：")
    for row in table_data:
        print(f"  {row}")

    if not table_data:
        print("错误：未解析到任何装修数据，请检查输入文件内容。")
        return {"status": "error", "message": "未解析到任何装修数据"}

    try:
        if os.path.exists(pptx_path):
            prs = Presentation(pptx_path)
        else:
            prs = Presentation()
    except Exception as e:
        print(f"加载 PowerPoint 时出错：{e}")
        return {"status": "error", "message": f"加载 PowerPoint 时出错：{e}"}
    
    while len(prs.slides) < 4:
        prs.slides.add_slide(prs.slide_layouts[6])
    
    slide = prs.slides[3]
    
    # 调整表格位置和尺寸，确保靠右下角
    rows = len(table_data) + 1
    cols = len(headers)
    left = Inches(8.5)  # 靠右，距左侧 5.5 英寸（页面宽 10 英寸 - 表格宽 4.0 英寸 - 右边距 0.5 英寸）
    top = Inches(3.0)   # 距顶部 3.0 英寸，保持在下半部分
    width = Inches(4.0) # 宽度占页面约 40%
    height = Inches(0.5 * rows)  # 每行 0.5 英寸，确保内容可见
    
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # 设置列宽和行高
    col_width = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_width
    for i in range(rows):
        table.rows[i].height = Inches(0.5)
    
    # 写入表头
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # 写入数据
    for row_idx, row_data in enumerate(table_data, 1):
        for col_idx, cell_data in enumerate(row_data):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_data
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 240, 240) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
    
    try:
        prs.save(pptx_path)
        print(f"表格成功插入到 {pptx_path} 的第 4 页（右下角，靠右）。")
        return {"status": "success", "message": f"表格成功插入到 {pptx_path} 的第 4 页"}
    except Exception as e:
        print(f"保存 PowerPoint 时出错：{e}")
        return {"status": "error", "message": f"保存 PowerPoint 时出错：{e}"}

if __name__ == "__main__":
    project_name = "华发四季半岛"
    result = run(project_name=project_name)
    print(result)