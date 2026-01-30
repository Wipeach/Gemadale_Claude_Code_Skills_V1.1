#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Insert opening information table from an Excel file into the second slide of a PowerPoint presentation with auto-adjusted column widths.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import pandas as pd
import os
from datetime import datetime

def run(project_name):
    """
    Insert opening information table from an Excel file into the second slide of a PowerPoint presentation.

    Args:
        project_name (str): Name of the project.

    Returns:
        dict: Contains status and output file path or error message.
    """
    timestamp = datetime.now().strftime("%Y%m%d")
    excel_file = f"resources/working_data/{project_name}_{timestamp}/processed_data/{project_name}_开盘信息.xlsx"
    pptx_path = f"resources/working_data/{project_name}_{timestamp}/processed_data/{project_name}_gemdale_housing_project_template.pptx"
    
    # Read Excel file
    try:
        df = pd.read_excel(excel_file)
        print(f"成功读取 Excel 文件：{excel_file}")
    except FileNotFoundError:
        print(f"错误：Excel 文件 {excel_file} 未找到。")
        return {"status": "error", "message": f"Excel 文件 {excel_file} 未找到"}
    except Exception as e:
        print(f"读取 Excel 文件时出错：{e}")
        return {"status": "error", "message": f"读取 Excel 文件时出错：{e}"}
    
    # Load PowerPoint
    try:
        if os.path.exists(pptx_path):
            prs = Presentation(pptx_path)
        else:
            print(f"错误：PowerPoint 文件 {pptx_path} 未找到。")
            return {"status": "error", "message": f"PowerPoint 文件 {pptx_path} 未找到"}
    except Exception as e:
        print(f"加载 PowerPoint 时出错：{e}")
        return {"status": "error", "message": f"加载 PowerPoint 时出错：{e}"}
    
    # Ensure at least 2 slides
    while len(prs.slides) < 2:
        prs.slides.add_slide(prs.slide_layouts[6])
    
    slide = prs.slides[1]  # Second slide (index 1)
    
    # Prepare table data
    headers = df.columns.tolist()
    table_data = df.values.tolist()
    rows = len(table_data) + 1  # Include header row
    cols = len(headers)
    
    # Table dimensions and position
    # Slide size is 13.33" x 7.5" (16:9)
    slide_width = prs.slide_width
    initial_table_width_in = 11.0  # maximum target width in inches
    # set row height (reduced a bit since font is smaller)
    row_height_in = 0.42  # inches per row
    table_height = Inches(row_height_in * rows)
    top = Inches(5.0)  # Slightly below midpoint
    # We'll compute left after we get total table width from content
    
    # Font sizes (reduced from prior values because你觉得字太大)
    header_font_pt = 11  # was 14
    data_font_pt = 9     # was 12

    # Helper: effective character length (Chinese-wide chars count double)
    def effective_char_count(s: str) -> int:
        if s is None:
            return 0
        s = str(s)
        cnt = 0
        for ch in s:
            # treat non-ASCII as wider (approx Chinese)
            cnt += 2 if ord(ch) > 127 else 1
        return cnt

    # Estimate per-column width in inches based on max effective chars and font size
    col_widths_in = []
    # Tunable coefficients:
    char_width_coef = 0.55  # factor to convert (char * font_pt) -> inches (empirical)
    padding_in = 0.18       # per-column padding (inches)
    min_col_width = 0.6     # minimum column width in inches
    max_col_width = initial_table_width_in * 0.7  # avoid a single column taking too much

    for col_idx in range(cols):
        # header effective length
        header = headers[col_idx]
        header_eff = effective_char_count(header)
        # data effective length (max among cells)
        max_data_eff = 0
        for row in table_data:
            # protect if row shorter
            val = row[col_idx] if col_idx < len(row) else ""
            max_data_eff = max(max_data_eff, effective_char_count(val))
        # compute two estimates (header and data) and take larger with corresponding font size
        width_from_header = (header_eff * header_font_pt * char_width_coef) / 72.0
        width_from_data = (max_data_eff * data_font_pt * char_width_coef) / 72.0
        est_w = max(width_from_header, width_from_data) + padding_in
        # clamp
        est_w = max(min_col_width, min(est_w, max_col_width))
        col_widths_in.append(est_w)

    total_calculated = sum(col_widths_in)
    # If too wide, scale down proportionally to initial_table_width_in
    if total_calculated > initial_table_width_in:
        scale_factor = initial_table_width_in / total_calculated
        col_widths_in = [w * scale_factor for w in col_widths_in]
        table_width_in = initial_table_width_in
    else:
        # Use the calculated total width so table is snug to content
        table_width_in = total_calculated

    # Convert to pptx lengths
    table_width = Inches(table_width_in)
    left = (slide_width - table_width) / 2  # center horizontally

    # Add table
    table = slide.shapes.add_table(rows, cols, left, top, table_width, table_height).table

    # Set column widths
    for i, width_in in enumerate(col_widths_in):
        try:
            table.columns[i].width = Inches(width_in)
        except Exception:
            # fallback to int EMU
            try:
                table.columns[i].width = int(Inches(width_in))
            except Exception:
                pass

    # Set row heights
    for i in range(rows):
        try:
            table.rows[i].height = Inches(row_height_in)
        except Exception:
            try:
                table.rows[i].height = int(Inches(row_height_in))
            except Exception:
                pass

    # Write headers
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = "" if header is None else str(header)
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(header_font_pt)
        p.font.bold = True
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
        p.font.color.rgb = RGBColor(255, 255, 255)  # White text
        # keep no wrap so width controls line length
        cell.text_frame.word_wrap = False

    # Write data
    for row_idx, row_data in enumerate(table_data, 1):
        for col_idx in range(cols):
            cell = table.cell(row_idx, col_idx)
            # guard missing columns
            try:
                cell_value = row_data[col_idx]
            except Exception:
                cell_value = ""
            cell.text = "" if cell_value is None else str(cell_value)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(data_font_pt)
            # zebra fill
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(240, 240, 240) if row_idx % 2 == 0 else RGBColor(255, 255, 255)
            cell.text_frame.word_wrap = False

    # Save PowerPoint
    try:
        prs.save(pptx_path)
        print(f"表格成功插入到 {pptx_path} 的第 2 页（中间靠下）。")
        return {"status": "success", "message": f"表格成功插入到 {pptx_path} 的第 2 页"}
    except Exception as e:
        print(f"保存 PowerPoint 时出错：{e}")
        return {"status": "error", "message": f"保存 PowerPoint 时出错：{e}"}

if __name__ == "__main__":
    result = run("华发四季半岛")
    print(result)
