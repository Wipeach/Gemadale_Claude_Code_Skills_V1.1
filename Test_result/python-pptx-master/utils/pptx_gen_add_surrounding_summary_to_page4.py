#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Insert surrounding information summary from a text file into a PowerPoint slide (page 4) at the top.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
import os
from datetime import datetime

def run(project_name):
    """
    读取已生成的周边信息总结，插入到 PowerPoint 第四页上方。
    
    Args:
        project_name (str): 项目名称。
    
    Returns:
        dict: {"status": "success/error", "message": "..."}
    """
    timestamp = datetime.now().strftime("%Y%m%d")
    input_file = f"resources/working_data/{project_name}_{timestamp}/processed_data/{project_name}_llm_周边信息.txt"
    pptx_path = f"resources/working_data/{project_name}_{timestamp}/processed_data/{project_name}_gemdale_housing_project_template.pptx"
    
    try:
        with open(input_file, 'r', encoding='utf-8') as f:
            surrounding_summary = f.read()
        print(f"成功读取周边信息总结文件：{input_file}")
        print("周边信息总结内容：")
        print(surrounding_summary)
    except FileNotFoundError:
        print(f"错误：文件 {input_file} 未找到。")
        return {"status": "error", "message": f"文件 {input_file} 未找到"}
    except Exception as e:
        print(f"读取文件时出错：{e}")
        return {"status": "error", "message": f"读取文件时出错：{e}"}
    
    try:
        if os.path.exists(pptx_path):
            prs = Presentation(pptx_path)
        else:
            print(f"错误：PowerPoint 文件 {pptx_path} 未找到。")
            return {"status": "error", "message": f"PowerPoint 文件 {pptx_path} 未找到"}
    except Exception as e:
        print(f"加载 PowerPoint 时出错：{e}")
        return {"status": "error", "message": f"加载 PowerPoint 时出错：{e}"}
    
    while len(prs.slides) < 4:
        prs.slides.add_slide(prs.slide_layouts[6])
    
    slide = prs.slides[3]
    
    left = Inches(0.5)
    top = Inches(0.5)
    width = Inches(9.0)
    height = Inches(3.5)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = "周边配套信息总结："
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.bold = True
    
    for line in surrounding_summary.split("\n"):
        if line.strip():
            p = text_frame.add_paragraph()
            p.text = line.strip()
            p.font.size = Pt(12)
            if line.startswith("- ") or line.startswith("* "):
                p.level = 1
            else:
                p.level = 0
    
    try:
        prs.save(pptx_path)
        print(f"周边信息总结成功插入到 {pptx_path} 的第 4 页上方。")
        return {"status": "success", "message": f"周边信息总结成功插入到 {pptx_path} 的第 4 页上方"}
    except Exception as e:
        print(f"保存 PowerPoint 时出错：{e}")
        return {"status": "error", "message": f"保存 PowerPoint 时出错：{e}"}

if __name__ == "__main__":
    project_name = "华发四季半岛"
    result = run(project_name=project_name)
    print(result)